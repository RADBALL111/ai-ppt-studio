#!/usr/bin/env python3
"""
make_video.py — AI 短视频全自动生成器  v9
一键：选题 → 分镜 → 多图源搜图 → TTS配音 → HTML+PPTX+MP4三出 → BGM混音 → 字幕

v9 新增: 字幕烧录 / Few-shot分镜示例 / 智能配色匹配
"""
import os, json, re, random, asyncio, shutil, subprocess, sqlite3, time, hashlib, warnings
import requests, edge_tts, numpy as np

# 抑制 moviepy FFmpeg 警告
warnings.filterwarnings("ignore", message=".*FFMPEG.*")
warnings.filterwarnings("ignore", message=".*proc.*")
from anthropic import Anthropic

# ==================== 配置 ====================
# DeepSeek
try:
    from config import DS_KEY as DS_API_KEY, PEXELS_KEY
except ImportError:
    DS_API_KEY = "sk-你的Key"
    PEXELS_KEY = "你的Key"
DS_BASE = "https://api.deepseek.com"
# Pexels
# Unsplash (free tier: 50 req/h)
UNSPLASH_KEY = None  # 可选，不填只用 Pexels

OUTPUT_DIR = "output"
CACHE_DB = os.path.join(OUTPUT_DIR, "image_cache.db")
W, H = 1920, 1080
FPS, BITRATE = 30, "6000k"
BGM_PATH = "bgm.mp3"
FONT_PATH = "/System/Library/Fonts/STHeiti Light.ttc"

os.makedirs(OUTPUT_DIR, exist_ok=True)

# ==================== 55套设计系统（融合四大开源项目精华） ====================
DESIGN_SYSTEMS = {
    # ===== 瑞士网格系（歸藏 guizang）=====
    "swiss_ikb":{"--bg":"#f8f8f8","--card":"#fff","--accent":"#0015FF","--accent2":"#000","--text":"#1a1a1a","--textDim":"#666","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "swiss_dark":{"--bg":"#1a1a1a","--card":"#2a2a2a","--accent":"#FFE600","--accent2":"#fff","--text":"#f0f0f0","--textDim":"#999","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "swiss_mint":{"--bg":"#f5faf7","--card":"#fff","--accent":"#00FF66","--accent2":"#0015FF","--text":"#111","--textDim":"#555","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "swiss_orange":{"--bg":"#fafaf8","--card":"#fff","--accent":"#FF3300","--accent2":"#000","--text":"#1a1a1a","--textDim":"#777","--fontTitle":"900","--radius":"0px","--grid":"swiss"},

    # ===== Light & Calm 浅色/平静（html-ppt-skill）=====
    "minimal_white":{"--bg":"#ffffff","--card":"#f5f5f5","--accent":"#111111","--accent2":"#666666","--text":"#111111","--textDim":"#888888","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "editorial_serif":{"--bg":"#fcfaf7","--card":"#ffffff","--accent":"#8b0000","--accent2":"#333333","--text":"#1a1a1a","--textDim":"#777777","--fontTitle":"900","--radius":"0px","--grid":"editorial"},
    "soft_pastel":{"--bg":"#fdf6f0","--card":"#ffffff","--accent":"#e8a0bf","--accent2":"#ba94b8","--text":"#4a3a4a","--textDim":"#9a8a9a","--fontTitle":"700","--radius":"12px","--grid":"none"},
    "solarized_light":{"--bg":"#fdf6e3","--card":"#eee8d5","--accent":"#268bd2","--accent2":"#2aa198","--text":"#073642","--textDim":"#586e75","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "catppuccin_latte":{"--bg":"#eff1f5","--card":"#e6e9ef","--accent":"#1e66f5","--accent2":"#fe640b","--text":"#4c4f69","--textDim":"#9ca0b0","--fontTitle":"800","--radius":"8px","--grid":"none"},
    "xiaohongshu":{"--bg":"#ffffff","--card":"#f8f8f8","--accent":"#ff2442","--accent2":"#cc1030","--text":"#333333","--textDim":"#999999","--fontTitle":"800","--radius":"8px","--grid":"none"},
    "paper_ink":{"--bg":"#faf8f5","--card":"#fff","--accent":"#2c2c2c","--accent2":"#666","--text":"#111","--textDim":"#777","--fontTitle":"800","--radius":"0px","--grid":"none"},
    "apple_snow":{"--bg":"#f5f5f7","--card":"#fff","--accent":"#007aff","--accent2":"#5856d6","--text":"#1d1d1f","--textDim":"#86868b","--fontTitle":"700","--radius":"12px","--grid":"none"},
    "warm_craft":{"--bg":"#faf5eb","--card":"#fff","--accent":"#d97706","--accent2":"#92400e","--text":"#44403c","--textDim":"#78716c","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "cream_paper":{"--bg":"#fefdf9","--card":"#fffef8","--accent":"#b45309","--accent2":"#78350f","--text":"#292524","--textDim":"#78716c","--fontTitle":"800","--radius":"2px","--grid":"none"},
    "japanese_minimal":{"--bg":"#f7f4ef","--card":"#ffffff","--accent":"#4a3728","--accent2":"#8b7355","--text":"#2c2420","--textDim":"#8b8682","--fontTitle":"700","--radius":"0px","--grid":"none"},

    # ===== Bold & Statement 大胆/宣言 =====
    "sharp_mono":{"--bg":"#ffffff","--card":"#000000","--accent":"#000000","--accent2":"#ff0000","--text":"#000000","--textDim":"#666666","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "neo_brutalism":{"--bg":"#ffffff","--card":"#ffdd00","--accent":"#ff0000","--accent2":"#0000ff","--text":"#000000","--textDim":"#555555","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "bauhaus":{"--bg":"#f0f0e8","--card":"#e8c547","--accent":"#c82828","--accent2":"#2c4c9c","--text":"#1a1a1a","--textDim":"#555555","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "swiss_grid":{"--bg":"#fafafa","--card":"#ffffff","--accent":"#ff3366","--accent2":"#333333","--text":"#111111","--textDim":"#777777","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "memphis_pop":{"--bg":"#fdf0e8","--card":"#ffffff","--accent":"#ff6b6b","--accent2":"#4ecdc4","--text":"#2d2d2d","--textDim":"#888888","--fontTitle":"900","--radius":"12px","--grid":"none"},
    "bold_signal":{"--bg":"#000","--card":"#0f0f0f","--accent":"#FF4500","--accent2":"#FFD700","--text":"#fff","--textDim":"#aaa","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "editorial_bw":{"--bg":"#fff","--card":"#fafafa","--accent":"#111","--accent2":"#e60000","--text":"#000","--textDim":"#666","--fontTitle":"900","--radius":"0px","--grid":"editorial"},

    # ===== Cool & Dark 冷暗色 =====
    "catppuccin_mocha":{"--bg":"#1e1e2e","--card":"#313244","--accent":"#cba6f7","--accent2":"#f5c2e7","--text":"#cdd6f4","--textDim":"#a6adc8","--fontTitle":"900","--radius":"6px","--grid":"none"},
    "dracula":{"--bg":"#282a36","--card":"#383a59","--accent":"#bd93f9","--accent2":"#ff79c6","--text":"#f8f8f2","--textDim":"#6272a4","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "tokyo_night":{"--bg":"#1a1b2e","--card":"#252740","--accent":"#ff9e64","--accent2":"#e0894f","--text":"#c0caf5","--textDim":"#6c7096","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "nord":{"--bg":"#2e3440","--card":"#3b4252","--accent":"#88c0d0","--accent2":"#81a1c1","--text":"#eceff4","--textDim":"#d8dee9","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "gruvbox_dark":{"--bg":"#282828","--card":"#3c3836","--accent":"#fabd2f","--accent2":"#fe8019","--text":"#ebdbb2","--textDim":"#a89984","--fontTitle":"900","--radius":"4px","--grid":"none"},
    "rose_pine":{"--bg":"#191724","--card":"#1f1d2e","--accent":"#ebbcba","--accent2":"#f6c177","--text":"#e0def4","--textDim":"#908caa","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "arctic_cool":{"--bg":"#0f1b2d","--card":"#152238","--accent":"#66ccff","--accent2":"#99ddff","--text":"#e6f3ff","--textDim":"#99bbdd","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "cyberpunk":{"--bg":"#0a0a1a","--card":"#1a1a3e","--accent":"#00ff88","--accent2":"#00cc66","--text":"#e0e0ff","--textDim":"#a0a0cc","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "pitch_navy":{"--bg":"#020617","--card":"#0f172a","--accent":"#38bdf8","--accent2":"#0284c7","--text":"#f8fafc","--textDim":"#94a3b8","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "obsidian":{"--bg":"#09090b","--card":"#18181b","--accent":"#fafafa","--accent2":"#a1a1aa","--text":"#fafafa","--textDim":"#a1a1aa","--fontTitle":"700","--radius":"2px","--grid":"none"},
    "noir_cinema":{"--bg":"#0a0a0a","--card":"#171717","--accent":"#e5e5e5","--accent2":"#737373","--text":"#f5f5f5","--textDim":"#a3a3a3","--fontTitle":"900","--radius":"0px","--grid":"none"},

    # ===== Warm & Vibrant 温暖/活力 =====
    "sunset_warm":{"--bg":"#1a0a05","--card":"#3a1a10","--accent":"#f97316","--accent2":"#ef4444","--text":"#fff7ed","--textDim":"#fdba74","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "coral_reef":{"--bg":"#1a0808","--card":"#2a1010","--accent":"#ff6b6b","--accent2":"#ffd93d","--text":"#fff0f0","--textDim":"#f0b0b0","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "desert_dusk":{"--bg":"#1c1008","--card":"#2a1a10","--accent":"#fb923c","--accent2":"#fbbf24","--text":"#fff7ed","--textDim":"#fed7aa","--fontTitle":"900","--radius":"4px","--grid":"none"},
    "vintage_film":{"--bg":"#1a1410","--card":"#2a2018","--accent":"#d4a574","--accent2":"#8b6914","--text":"#f5e6d3","--textDim":"#c4a882","--fontTitle":"800","--radius":"2px","--grid":"none"},

    # ===== Effect-Heavy 特效系 =====
    "glassmorphism":{"--bg":"#0f0f23","--card":"rgba(255,255,255,0.08)","--accent":"#ffffff","--accent2":"#c0c0ff","--text":"#ffffff","--textDim":"#c0c0d0","--fontTitle":"800","--radius":"16px","--grid":"none"},
    "aurora":{"--bg":"#0a0a1a","--card":"#0f1030","--accent":"#00ffff","--accent2":"#ff00ff","--text":"#e0ffff","--textDim":"#80c0c0","--fontTitle":"900","--radius":"12px","--grid":"none"},
    "terminal_green":{"--bg":"#0c0c0c","--card":"#1a1a1a","--accent":"#00ff00","--accent2":"#00cc00","--text":"#ccffcc","--textDim":"#66aa66","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "blueprint":{"--bg":"#0a1a3a","--card":"#0f2850","--accent":"#ffffff","--accent2":"#88ccff","--text":"#ffffff","--textDim":"#88aacc","--fontTitle":"800","--radius":"2px","--grid":"none"},
    "neon_synthwave":{"--bg":"#0f0b1a","--card":"#1a1030","--accent":"#ff00ff","--accent2":"#00ffff","--text":"#f0e0ff","--textDim":"#c080e0","--fontTitle":"900","--radius":"12px","--grid":"none"},
    "mono_void":{"--bg":"#000","--card":"#0a0a0a","--accent":"#fff","--accent2":"#808080","--text":"#fff","--textDim":"#808080","--fontTitle":"900","--radius":"0px","--grid":"none"},

    # ===== v2 Professional 专业版 =====
    "corporate_clean":{"--bg":"#f8fafc","--card":"#ffffff","--accent":"#2563eb","--accent2":"#1d4ed8","--text":"#0f172a","--textDim":"#64748b","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "pitch_deck_vc":{"--bg":"#0c0c1d","--card":"#162040","--accent":"#818cf8","--accent2":"#6366f1","--text":"#e0e7ff","--textDim":"#a5b4fc","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "academic_paper":{"--bg":"#ffffff","--card":"#fafafa","--accent":"#1e40af","--accent2":"#3b82f6","--text":"#1e293b","--textDim":"#64748b","--fontTitle":"700","--radius":"0px","--grid":"none"},
    "engineering_wp":{"--bg":"#f0f4f8","--card":"#ffffff","--accent":"#0284c7","--accent2":"#0369a1","--text":"#1a202c","--textDim":"#718096","--fontTitle":"800","--radius":"2px","--grid":"none"},
    "corp_blue":{"--bg":"#0a1e3d","--card":"#112850","--accent":"#3b82f6","--accent2":"#1d4ed8","--text":"#dbeafe","--textDim":"#93c5fd","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "emerald_deck":{"--bg":"#022c22","--card":"#064e3b","--accent":"#34d399","--accent2":"#6ee7b7","--text":"#ecfdf5","--textDim":"#a7f3d0","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "slate_pro":{"--bg":"#1e293b","--card":"#334155","--accent":"#e2e8f0","--accent2":"#94a3b8","--text":"#f8fafc","--textDim":"#cbd5e1","--fontTitle":"800","--radius":"4px","--grid":"none"},

    # ===== v2 Editorial 杂志版 =====
    "magazine_bold":{"--bg":"#ffffff","--card":"#f5f5f5","--accent":"#dc2626","--accent2":"#111111","--text":"#111111","--textDim":"#666666","--fontTitle":"900","--radius":"0px","--grid":"editorial"},
    "news_broadcast":{"--bg":"#0a0a10","--card":"#1a1a28","--accent":"#ff3333","--accent2":"#ffffff","--text":"#ffffff","--textDim":"#999999","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "midcentury":{"--bg":"#fdf6e3","--card":"#ffffff","--accent":"#e07a5f","--accent2":"#3d405b","--text":"#2d2d2d","--textDim":"#777777","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "retro_tv":{"--bg":"#111118","--card":"#1a1a28","--accent":"#00ffcc","--accent2":"#ff6699","--text":"#e0e0e0","--textDim":"#888888","--fontTitle":"900","--radius":"12px","--grid":"none"},

    # ===== v2 Dramatic 戏剧性 =====
    "cyberpunk_neon":{"--bg":"#08001a","--card":"#130030","--accent":"#ff00ff","--accent2":"#00ffff","--text":"#f0e0ff","--textDim":"#b090e0","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "vaporwave":{"--bg":"#0d0d1a","--card":"#1a1a3e","--accent":"#ff6ac1","--accent2":"#00ffff","--text":"#ffe0f0","--textDim":"#a080c0","--fontTitle":"900","--radius":"16px","--grid":"none"},
    "y2k_chrome":{"--bg":"#e8e8f0","--card":"#ffffff","--accent":"#ff1493","--accent2":"#00bfff","--text":"#1a1a2e","--textDim":"#666688","--fontTitle":"900","--radius":"20px","--grid":"none"},

    # ===== 自然系 =====
    "forest_moss":{"--bg":"#0a1a10","--card":"#142818","--accent":"#4ade80","--accent2":"#22c55e","--text":"#dcfce7","--textDim":"#86efac","--fontTitle":"900","--radius":"6px","--grid":"none"},
    "ocean_depth":{"--bg":"#051a2a","--card":"#0a2a40","--accent":"#06b6d4","--accent2":"#0891b2","--text":"#ecfeff","--textDim":"#67e8f9","--fontTitle":"800","--radius":"8px","--grid":"none"},
    "arctic_frost":{"--bg":"#f0f9ff","--card":"#fff","--accent":"#0ea5e9","--accent2":"#6366f1","--text":"#0f172a","--textDim":"#64748b","--fontTitle":"700","--radius":"8px","--grid":"none"},
    "electric_lime":{"--bg":"#0a0f00","--card":"#141f00","--accent":"#a3ff00","--accent2":"#00ff88","--text":"#f0ffe0","--textDim":"#90b060","--fontTitle":"900","--radius":"4px","--grid":"none"},
    "brutalist":{"--bg":"#fff","--card":"#000","--accent":"#f00","--accent2":"#00f","--text":"#000","--textDim":"#666","--fontTitle":"900","--radius":"0px","--grid":"none"},
}

print(f"🎨 已加载 {len(DESIGN_SYSTEMS)} 套设计系统", flush=True)

# ==================== SQLite 图搜缓存 ====================
def init_cache():
    conn = sqlite3.connect(CACHE_DB)
    conn.execute("CREATE TABLE IF NOT EXISTS images(keyword TEXT PRIMARY KEY, path TEXT, photographer TEXT, ts REAL)")
    conn.commit(); return conn

def cache_get(conn, keyword):
    row = conn.execute("SELECT path,photographer FROM images WHERE keyword=?",(keyword,)).fetchone()
    return row if row else (None,None)

def cache_set(conn, keyword, path, photographer):
    conn.execute("INSERT OR REPLACE INTO images VALUES(?,?,?,?)",(keyword,path,photographer,time.time()))
    conn.commit()

# ==================== P1: 多图源搜图 ====================
def search_pexels(keyword):
    """Pexels + locale=zh-CN"""
    headers = {"Authorization": PEXELS_KEY}
    has_cjk = any('一' <= c <= '鿿' for c in keyword)
    params = {"query": keyword, "per_page": 1, "orientation": "landscape"}
    if has_cjk: params["locale"] = "zh-CN"
    try:
        r = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params, timeout=10)
        data = r.json()
        if "photos" in data and data["photos"]:
            p = data["photos"][0]; return p["src"]["large"], p["photographer"]
    except: pass
    return None, None

def search_unsplash(keyword):
    """Unsplash free tier fallback"""
    if not UNSPLASH_KEY: return None, None
    try:
        r = requests.get("https://api.unsplash.com/search/photos",
                         headers={"Authorization": f"Client-ID {UNSPLASH_KEY}"},
                         params={"query": keyword, "per_page": 1, "orientation": "landscape"}, timeout=10)
        data = r.json()
        if "results" in data and data["results"]:
            p = data["results"][0]; return p["urls"]["regular"], p["user"]["name"]
    except: pass
    return None, None

def fetch_image(keyword, index):
    """多源搜图 + 缓存"""
    conn = init_cache()
    # 1. 缓存命中
    cached_path, cached_photog = cache_get(conn, keyword)
    if cached_path and os.path.exists(cached_path):
        return cached_path, cached_photog + " (缓存)"
    # 2. Pexels
    img_dir = os.path.join(OUTPUT_DIR, "images")
    os.makedirs(img_dir, exist_ok=True)
    url, photog = search_pexels(keyword)
    # 3. Unsplash fallback
    if not url: url, photog = search_unsplash(keyword)
    if not url: return None, None
    # 4. 下载
    try:
        data = requests.get(url, timeout=15).content
        path = os.path.join(img_dir, f"img_{index}.jpg")
        with open(path, "wb") as f: f.write(data)
        cache_set(conn, keyword, path, photog)
        return path, photog
    except: return None, None

# ==================== P2: Structured Output (DeepSeek原生JSON) ====================
def ask_structured(system_prompt, user_prompt, max_tokens=3000):
    """DeepSeek 原生 API JSON mode"""
    headers = {"Authorization": f"Bearer {DS_API_KEY}", "Content-Type": "application/json"}
    payload = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "max_tokens": max_tokens,
        "response_format": {"type": "json_object"},
        "temperature": 0.7,
    }
    for attempt in range(3):
        try:
            r = requests.post(f"{DS_BASE}/v1/chat/completions", headers=headers, json=payload, timeout=60)
            data = r.json()
            content = data["choices"][0]["message"]["content"]
            return json.loads(content)
        except Exception as e:
            if attempt == 2: raise
            time.sleep(1)
    return {}

def ask_text(prompt, max_tokens=500):
    """简单文本请求"""
    headers = {"Authorization": f"Bearer {DS_API_KEY}", "Content-Type": "application/json"}
    payload = {"model": "deepseek-chat", "messages": [{"role": "user", "content": prompt}],
               "max_tokens": max_tokens}
    r = requests.post(f"{DS_BASE}/v1/chat/completions", headers=headers, json=payload, timeout=30)
    return r.json()["choices"][0]["message"]["content"].strip()

# ==================== AI 实时主题生成（无限模板） ====================

THEME_GEN_PROMPT = """你是顶级视觉设计师。根据选题生成一套专属配色方案。

## 要求
- 输出 7 个 CSS 变量：--bg(背景) --card(卡片) --accent(强调色) --accent2(辅助强调) --text(主文字) --textDim(次级文字)
- 所有颜色用 #RRGGBB hex 格式
- --fontTitle: 700/800/900 三选一
- --radius: 0px/4px/8px/12px/16px 五选一
- --grid: "swiss"/"editorial"/"none" 三选一
- 色彩要有情感倾向：选题是负面情绪用冷暗色、正面用暖亮色、震惊用高对比撞色
- 不要紫色渐变、不要超过 3 个主色

输出纯JSON：{"--bg":"#...","--card":"#...","--accent":"#...","--accent2":"#...","--text":"#...","--textDim":"#...","--fontTitle":"900","--radius":"4px","--grid":"none","name":"主题名"}"""

def generate_ai_theme(topic):
    """AI 根据选题实时生成专属主题"""
    try:
        result = ask_structured(THEME_GEN_PROMPT, f"选题：{topic}", 1000)
        # 确保有必需的键
        required = ["--bg","--accent","--text"]
        if all(k in result for k in required):
            name = result.pop("name", topic[:8])
            print(f"  🤖 AI生成主题: {name}")
            return result
    except: pass
    return None


def pick_theme(topic):
    """优先 AI 实时生成，失败则智能匹配 62 套库，再失败随机"""
    # 1. AI 实时生成（50% 概率）
    if random.random() < 0.5:
        ai_theme = generate_ai_theme(topic)
        if ai_theme:
            return ai_theme, "AI实时生成"
    # 2. 智能匹配
    theme_name = smart_theme(topic)
    if theme_name in DESIGN_SYSTEMS:
        return DESIGN_SYSTEMS[theme_name], theme_name
    # 3. 随机
    name = random.choice(list(DESIGN_SYSTEMS.keys()))
    return DESIGN_SYSTEMS[name], name
THEME_KEYWORDS = {
    # 瑞士系
    "swiss_ikb":       ["科技","AI","数据","算法","代码","程序","黑客","数字","软件","工程师"],
    "swiss_dark":      ["批判","真相","揭露","骗局","陷阱","暗黑","反常识","打脸"],
    "swiss_mint":      ["环保","绿色","自然","植物","气候","地球","生态","可持续"],
    "swiss_orange":    ["安全","警告","危险","紧急","消防","救援"],
    # 浅色/平静
    "minimal_white":   ["极简","纯粹","本质","干净","留白","空白"],
    "editorial_serif": ["杂志","编辑","文艺","专栏","深度","长篇","报道"],
    "soft_pastel":     ["温柔","少女","治愈","甜蜜","可爱","萌"],
    "solarized_light": ["编程","代码","终端","开发","技术文档","开源"],
    "catppuccin_latte":["日常","笔记","日记","轻松","咖啡","慢生活"],
    "xiaohongshu":     ["小红书","种草","穿搭","美妆","好物","分享"],
    "paper_ink":       ["读书","写作","思考","哲学","人文","墨水","纸"],
    "apple_snow":      ["产品","设计","Apple","简洁","发布会","新品"],
    "japanese_minimal":["侘寂","日式","禅","茶道","简约","空"],
    # 大胆/宣言
    "sharp_mono":      ["对比","极端","冲击","黑白","对立","撕逼"],
    "neo_brutalism":   ["新粗野主义","反设计","大胆","年轻","街头","潮"],
    "bauhaus":         ["包豪斯","几何","现代主义","构成","抽象","艺术"],
    "swiss_grid":      ["网格","系统","秩序","排版","结构","精确"],
    "memphis_pop":     ["孟菲斯","波普","活泼","有趣","卡通","欢乐"],
    "bold_signal":     ["震撼","爆炸","头条","突发","危机","挑战"],
    "editorial_bw":    ["社论","评论","黑白","立场","态度","锋芒"],
    # 冷暗色
    "catppuccin_mocha":["深夜","暗色","柔和","舒适","极客","编码"],
    "dracula":         ["吸血鬼","暗黑","神秘","哥特","酷","炫"],
    "tokyo_night":     ["东京","夜景","霓虹","日本","赛博","城市"],
    "nord":            ["北欧","冰雪","干净","冷调","理性","极简"],
    "gruvbox_dark":    ["复古","温暖","暗色","怀旧","咖啡色"],
    "rose_pine":       ["玫瑰","优雅","暗色","柔和","低饱和"],
    "arctic_cool":     ["北极","冰冷","清爽","科技感","蓝调"],
    "cyberpunk":       ["赛博朋克","未来","虚拟","霓虹","黑客"],
    "pitch_navy":      ["商业","融资","创业","增长","市场","战略","Pitch"],
    "obsidian":        ["笔记","知识库","第二大脑","链接","图谱"],
    # 温暖/活力
    "sunset_warm":     ["落日","温暖","黄昏","夏天","热烈","激情"],
    "coral_reef":      ["珊瑚","海洋","活力","热情","夏日","海滩"],
    "vintage_film":    ["复古","怀旧","胶片","老照片","经典","电影"],
    # 特效系
    "glassmorphism":   ["玻璃","透明","毛玻璃","层次","现代","iOS"],
    "aurora":          ["极光","梦幻","炫彩","流光","魔法","幻境"],
    "terminal_green":  ["终端","黑客","命令行","复古电脑","DOS"],
    "blueprint":       ["蓝图","工程","建筑","计划","方案","架构"],
    "neon_synthwave":  ["合成波","80年代","霓虹","复古未来","电子"],
    "mono_void":       ["空洞","虚无","极夜","极致暗","沉浸"],
    # 专业
    "corporate_clean": ["企业","商务","专业","汇报","年报","正式"],
    "pitch_deck_vc":   ["VC","投资人","路演","Demo","融资","估值"],
    "academic_paper":  ["学术","论文","研究","大学","教授","科学"],
    "engineering_wp":  ["工程","白皮书","技术","架构","系统"],
    "corp_blue":       ["职场","管理","效率","方法","上班","公司"],
    "emerald_deck":    ["金融","财富","理财","投资","基金","银行"],
    # 杂志
    "magazine_bold":   ["头条","封面","杂志","冲击","大标题","新闻"],
    "news_broadcast":  ["新闻","电视台","直播","播报","快讯"],
    "midcentury":      ["中世纪","经典","复古50s","设计史","广告"],
    "retro_tv":        ["老电视","雪花屏","复古","VHS","怀旧电视"],
    # 戏剧
    "cyberpunk_neon":  ["赛博朋克","霓虹灯","未来城市","科幻","黑客帝国"],
    "vaporwave":       ["蒸汽波","80s","粉蓝","塑料","迷幻","网络"],
    "y2k_chrome":      ["千禧年","Y2K","金属","反光","未来感","镀铬"],
    # 自然
    "forest_moss":     ["森林","苔藓","绿色","户外","探险","野生"],
    "ocean_depth":     ["海洋","深海","水底","潜水","蓝色深渊"],
    "desert_dusk":     ["沙漠","黄昏","旷野","孤独","西部","旅行"],
    "warm_craft":      ["手作","咖啡","美食","面包","慢","工匠"],
    "cream_paper":     ["纸张","奶油","温暖","信纸","日记","手工"],
}

def smart_theme(topic):
    """根据选题关键词匹配最佳主题"""
    scores = {}
    for theme, keywords in THEME_KEYWORDS.items():
        scores[theme] = sum(1 for kw in keywords if kw in topic)
    best = max(scores, key=scores.get)
    if scores[best] > 0:
        return best
    # 无匹配：瑞士网格或随机
    return random.choice(["swiss_ikb","swiss_dark","paper_ink","apple_snow","bold_signal"])


# ==================== P5: 分镜生成 + Few-shot ====================
FEWSHOT_EXAMPLES = """
## 高质量分镜示例

### 示例1: "你的工资为什么永远涨不过物价"
[
  {"layout":"hero","rhythm":"hero","title":"你的工资涨了","subtitle":"但物价涨得更快","highlight":"工资陷阱"},
  {"layout":"kpi_bold","rhythm":"non-hero","title":"过去10年","subtitle":"物价翻3倍，工资涨1.5倍","highlight":"3x差距"},
  {"layout":"quote_slide","rhythm":"hero","title":"你不是在赚钱","subtitle":"你是在被通胀慢慢吃掉","highlight":"被吃掉"},
  {"layout":"split_col","rhythm":"non-hero","title":"你以为存银行安全","subtitle":"实际上每年贬值6%","highlight":"贬值"},
  {"layout":"hero","rhythm":"hero","title":"唯一的出路","subtitle":"让钱为你工作","highlight":"资产"},
]

### 示例2: "为什么内向的人更适合做领导"
[
  {"layout":"hero","rhythm":"hero","title":"你以为领导都要会社交","subtitle":"数据告诉你恰恰相反","highlight":"反常识"},
  {"layout":"statement","rhythm":"non-hero","title":"哈佛研究：内向者CEO比例高达40%","subtitle":"远超人群中25%的占比","highlight":"40%"},
  {"layout":"quote_slide","rhythm":"hero","title":"沉默不是弱点","subtitle":"是深度思考的外在表现","highlight":"沉默"},
  {"layout":"split_col","rhythm":"non-hero","title":"内向者三大优势","subtitle":"深度倾听 / 独立决策 / 不依赖外部认可","highlight":"优势"},
  {"layout":"hero","rhythm":"hero","title":"世界需要安静的力量","subtitle":"你不需要变成另一个人","highlight":"安静的"},
]

"""

STORYBOARD_PROMPT = """你是短视频导演。根据选题生成分镜脚本。

## 参考示例（学习节奏和冲突感）
{FEWSHOT}

## 绝对禁止
- 禁止鸡汤废话（"让我们一起""加油""相信自己""你是最棒的"）
- 禁止每句话超过15字
- 禁止连续两页同一布局
- 禁止平淡开头（第一页必须有数字/反常识/冲突）

## 强制规则
1. hero/non-hero 交替
2. 前3页必须建立冲突或悬念
3. highlight 是 2-6 字关键词
4. layout: hero / statement / split_col / full_bleed / quote_slide / kpi_bold

## 自评（不满足就重写）
- 第一页是否够炸？
- 有没有鸡汤废话？
- 有没有数字/数据支撑？

输出 6-10 页JSON：{{"storyboard":[...]}}"""

def generate_storyboard(topic):
    """分镜生成 + Few-shot + 自评重写（最多2轮）"""
    prompt = STORYBOARD_PROMPT.replace("{FEWSHOT}", FEWSHOT_EXAMPLES)
    for attempt in range(2):
        try:
            result = ask_structured(prompt, f"选题：{topic}", 4000)
            sb = result.get("storyboard", result if isinstance(result, list) else [])
            if not isinstance(sb, list) or len(sb) < 3:
                continue
            bad_words = ["我们一起","加油","相信","努力就能","你是最棒的"]
            bad = sum(1 for s in sb for w in bad_words if w in (s.get("title","")+s.get("subtitle","")))
            if bad > 0 and attempt == 0:
                continue
            return sb
        except: continue
    return [{"layout":"hero","rhythm":"hero","title":topic,"subtitle":"","highlight":"真相"},
            {"layout":"statement","rhythm":"non-hero","title":"颠覆认知","subtitle":"","highlight":"关键"}]

# ==================== TTS ====================
async def _tts(text, path):
    await edge_tts.Communicate(text, voice="zh-CN-YunxiNeural").save(path)

def gen_voice(storyboard):
    print("  🎙 生成配音...")
    async def b():
        await asyncio.gather(*[_tts(s.get("title",""),os.path.join(OUTPUT_DIR,f"voice_{i}.mp3")) for i,s in enumerate(storyboard)])
    asyncio.run(b())
    durs = []
    for i in range(len(storyboard)):
        try:
            from moviepy import AudioFileClip
            ac = AudioFileClip(os.path.join(OUTPUT_DIR,f"voice_{i}.mp3"))
            durs.append(ac.duration + 0.3); ac.close()
        except: durs.append(2.0)
    return durs

# ==================== P3: 自动氛围音轨 ====================
def generate_ambient(duration_sec, output_path):
    """pydub 生成简单氛围音轨（低频正弦波+白噪声）"""
    try:
        from pydub import AudioSegment
        from pydub.generators import Sine
        sample_rate = 44100
        # 低频垫音 C3 (~131Hz) + E3 (~165Hz)
        low = Sine(131).to_audio_segment(duration=duration_sec * 1000)
        mid = Sine(165).to_audio_segment(duration=duration_sec * 1000)
        ambient = low.overlay(mid)
        ambient = ambient - 25  # 很轻的背景
        ambient = ambient.fade_in(2000).fade_out(3000)
        ambient.export(output_path, format="mp3", bitrate="128k")
        return output_path
    except Exception as e:
        print(f"  ⚠️ 氛围音生成失败: {e}")
        return None

# ==================== 字幕生成 + 烧录 ====================

def generate_srt(storyboard, voice_durs):
    """从分镜+配音时长生成 SRT 字幕"""
    srt_lines = []
    t = 0.0
    for i, slide in enumerate(storyboard):
        dur = voice_durs[i] if i < len(voice_durs) else 2.0
        title = slide.get("title", "")
        subtitle = slide.get("subtitle", "")
        # 合并标题+副标题，限制长度避免溢出
        text = title
        if subtitle and len(text + " " + subtitle) < 25:
            text += "，" + subtitle
        if len(text) > 30:
            # 自动断行
            mid = len(text) // 2
            text = text[:mid] + "\n" + text[mid:]

        start = t; end = t + dur
        srt_lines.append(f"{i+1}")
        srt_lines.append(f"{_fmt_srt_time(start)} --> {_fmt_srt_time(end)}")
        srt_lines.append(text)
        srt_lines.append("")
        t += dur

    srt_path = os.path.join(OUTPUT_DIR, "subtitles.srt")
    with open(srt_path, "w", encoding="utf-8") as f:
        f.write("\n".join(srt_lines))
    return srt_path

def _fmt_srt_time(sec):
    h = int(sec // 3600); m = int((sec % 3600) // 60)
    s = int(sec % 60); ms = int((sec % 1) * 1000)
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"

def burn_subtitles(video_path, srt_path):
    """FFmpeg 烧录字幕，使用系统 STHeiti 中文字体"""
    out_path = video_path.replace(".mp4", "_sub.mp4")
    env = {"PATH": os.path.expanduser("~/bin") + ":" + os.environ.get("PATH", "")}
    # macOS 系统自带中文字体
    font_dir = "/System/Library/Fonts"
    font_name = "STHeiti"  # 华文黑体，macOS 自带，fontconfig 可识别
    try:
        r = subprocess.run([
            "ffmpeg", "-y", "-i", video_path,
            "-vf", f"subtitles={srt_path}:fontsdir={font_dir}:force_style='FontName={font_name},FontSize=28,PrimaryColour=&H00FFFFFF,OutlineColour=&H00000000,Outline=2,Shadow=1,MarginV=80'",
            "-c:v", "libx264", "-crf", "18", "-c:a", "copy",
            out_path
        ], capture_output=True, env=env, timeout=120)
        if os.path.exists(out_path) and os.path.getsize(out_path) > 1000:
            print("  📝 字幕已烧录")
            return out_path
        # fallback: 直接用字体文件路径
        font_file = "/System/Library/Fonts/STHeiti Light.ttc"
        r2 = subprocess.run([
            "ffmpeg", "-y", "-i", video_path,
            "-vf", f"subtitles={srt_path}:force_style='FontFile={font_file},FontSize=24,PrimaryColour=&H00FFFFFF,Outline=1,Shadow=1,MarginV=80'",
            "-c:v", "libx264", "-crf", "18", "-c:a", "copy",
            out_path
        ], capture_output=True, env=env, timeout=120)
        if os.path.exists(out_path) and os.path.getsize(out_path) > 1000:
            print("  📝 字幕已烧录（直接字体文件）")
            return out_path
    except Exception as e:
        print(f"  ⚠️ 字幕烧录失败: {e}")
    return video_path


# ==================== P3: BGM 混音(ducking) ====================
def mix_bgm(video_path, bgm_source, output_path):
    """配音 + BGM 混音，ducking -15dB"""
    try:
        from pydub import AudioSegment
        # 提取视频音频
        tmp_voice = os.path.join(OUTPUT_DIR, "_voice_tmp.wav")
        env = {"PATH": os.path.expanduser("~/bin") + ":" + os.environ.get("PATH", "")}
        subprocess.run(["ffmpeg","-y","-i",video_path,"-vn","-acodec","pcm_s16le",tmp_voice],
                       capture_output=True, env=env)
        if not os.path.exists(tmp_voice): return video_path

        voice = AudioSegment.from_wav(tmp_voice)
        # BGM: 文件或生成
        if os.path.exists(bgm_source):
            bgm = AudioSegment.from_mp3(bgm_source)
        elif bgm_source and os.path.exists(bgm_source):
            bgm = AudioSegment.from_mp3(bgm_source)
        else:
            # 自动生成
            gen_path = os.path.join(OUTPUT_DIR, "_ambient.mp3")
            if generate_ambient(len(voice) / 1000.0, gen_path):
                bgm = AudioSegment.from_mp3(gen_path)
            else:
                os.remove(tmp_voice); return video_path

        # 循环 BGM
        while len(bgm) < len(voice): bgm += bgm
        bgm = bgm[:len(voice)]
        bgm = bgm - 15  # ducking
        mixed = voice.overlay(bgm)

        tmp_mix = os.path.join(OUTPUT_DIR, "_mixed.aac")
        mixed.export(tmp_mix, format="adts")
        tmp_video = os.path.join(OUTPUT_DIR, "_tmp_video.mp4")
        subprocess.run(["ffmpeg","-y","-i",video_path,"-i",tmp_mix,
                        "-c:v","copy","-c:a","aac","-map","0:v:0","-map","1:a:0","-shortest",tmp_video],
                       capture_output=True, env=env)
        if os.path.exists(tmp_video):
            os.replace(tmp_video, output_path)
        os.remove(tmp_voice); os.remove(tmp_mix) if os.path.exists(tmp_mix) else None
        print("  🎵 BGM 混音完成")
        return output_path
    except Exception as e:
        print(f"  ⚠️ BGM 混音失败: {e}")
        return video_path

# ==================== P7: 增强 PPTX 导出 ====================
def export_pptx(storyboard, theme_data, images):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        ds = theme_data
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

        def hex_to_rgb(h):
            h = str(h).lstrip('#').split(',')[0].split(')')[0].strip()
            if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2  # #fff → ffffff
            if len(h) < 6: h = h.ljust(6, '0')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        bg_rgb = hex_to_rgb(ds.get("--bg","#000"))
        text_rgb = hex_to_rgb(ds.get("--text","#fff"))
        accent_rgb = hex_to_rgb(ds.get("--accent","#f00"))

        for i, slide_data in enumerate(storyboard):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = bg_rgb

            # 图片背景（有图时）
            if i in images and images[i][0]:
                img_path = images[i][0]
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                    # 半透明遮罩
                    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))  # rectangle
                    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0,0,0)
                    shape.fill.fore_color.brightness = -0.6
                    shape.line.fill.background()

            title = slide_data.get("title","")
            subtitle = slide_data.get("subtitle","")
            hl = slide_data.get("highlight","")

            if title:
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.5))
                tf = txBox.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = text_rgb

            if subtitle:
                txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.5))
                tf2 = txBox2.text_frame; tf2.word_wrap = True
                p2 = tf2.paragraphs[0]; p2.text = subtitle; p2.alignment = PP_ALIGN.CENTER
                p2.font.size = Pt(20); p2.font.color.rgb = text_rgb

            if hl:
                txBox3 = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2.3), Inches(0.6))
                tf3 = txBox3.text_frame
                p3 = tf3.paragraphs[0]; p3.text = hl; p3.alignment = PP_ALIGN.CENTER
                p3.font.size = Pt(18); p3.font.bold = True; p3.font.color.rgb = accent_rgb

        pptx_path = os.path.join(OUTPUT_DIR, "slides.pptx")
        prs.save(pptx_path)
        print(f"  📄 PPTX: {pptx_path}")
        return pptx_path
    except Exception as e:
        print(f"  ⚠️ PPTX: {e}")
        return None

# ==================== HTML 生成 ====================
def build_html(storyboard, theme_data, voice_durs, images):
    ds = theme_data  # 直接接受 CSS 变量字典（来自库或AI生成）
    is_swiss = ds.get("--grid") == "swiss"
    if images is None: images = {}

    slides, anis, audios, img_to_copy = [], [], [], []
    t = 0

    SWISS_LAYOUTS = {
        "hero": lambda t2,s,h: (f'<div class="swiss-hero"><div class="accent-bar"></div><h1>{t2}</h1><p class="lead">{s}</p><div class="tag">{h}</div></div>', 1.6),
        "statement":lambda t2,s,h: (f'<div class="swiss-statement"><div class="num">{h}</div><h2>{t2}</h2><p>{s}</p></div>', 1.0),
        "split_col":lambda t2,s,h: (f'<div class="swiss-split"><div class="col-l"><div class="label">{h}</div><h2>{t2}</h2></div><div class="col-r"><p>{s}</p></div></div>', 1.0),
        "quote_slide":lambda t2,s,h: (f'<div class="swiss-quote"><div class="rule"></div><blockquote>{t2}</blockquote><cite>{h}</cite></div>', 1.0),
        "kpi_bold":lambda t2,s,h: (f'<div class="swiss-kpi"><span class="num">{h}</span><span class="title">{t2}</span><span class="desc">{s}</span></div>', 1.0),
    }

    for i, slide in enumerate(storyboard):
        layout = slide.get("layout","statement")
        title = slide.get("title","")
        sub = slide.get("subtitle","")
        hl = slide.get("highlight","")
        dur = max(voice_durs[i], 1.8)
        if layout in ("hero","full_bleed","cover"): dur *= 1.6
        sid = f"s{i}"

        audios.append(f'<audio src="voice_{i}.mp3" data-start="{t:.1f}" data-duration="{dur:.1f}"></audio>')

        # 图片背景
        has_img = i in images and images[i][0]
        img_bg = ""
        if has_img:
            img_bg = f'<div class="img-bg" style="background-image:url(images/img_{i}.jpg)"></div><div class="img-overlay"></div>'
            img_to_copy.append((images[i][0], i))

        # 布局
        if is_swiss and layout in SWISS_LAYOUTS:
            body, _ = SWISS_LAYOUTS[layout](title, sub, hl)
            if layout=="hero":
                for sel,dy,dd in [(".accent-bar",0,0.6),("h1",0.2,0.7),(".lead",0.4,0.5),(".tag",0.5,0.4)]:
                    anis.append(f'tl.from("#{sid} {sel}",{{opacity:0,y:30,duration:{dd}}},{t+dy:.1f});')
            elif layout=="statement":
                for sel,dy,dd in [(".num",0,0.5),("h2",0.15,0.6),("p",0.3,0.5)]:
                    anis.append(f'tl.from("#{sid} {sel}",{{opacity:0,y:25,duration:{dd}}},{t+dy:.1f});')
            elif layout=="split_col":
                anis.append(f'tl.from("#{sid} .col-l",{{opacity:0,x:-40,duration:0.6}},{t:.1f});')
                anis.append(f'tl.from("#{sid} .col-r",{{opacity:0,x:40,duration:0.6}},{t+0.2:.1f});')
            elif layout=="quote_slide":
                anis.append(f'tl.from("#{sid} .rule",{{scaleX:0,duration:0.5}},{t:.1f});')
                anis.append(f'tl.from("#{sid} blockquote",{{opacity:0,y:30,duration:0.7}},{t+0.2:.1f});')
            elif layout=="kpi_bold":
                anis.append(f'tl.from("#{sid} .num",{{opacity:0,scale:0.6,duration:0.7,ease:"back.out(1.5)"}},{t:.1f});')
                anis.append(f'tl.from("#{sid} .title",{{opacity:0,y:20,duration:0.5}},{t+0.3:.1f});')
        else:
            body = f'<h2 id="{sid}_t">{title}</h2><p id="{sid}_s">{sub}</p><div id="{sid}_h" class="tag">{hl}</div>'
            for sel,dy,dd in [(f"#{sid}_t",0,0.6),(f"#{sid}_s",0.15,0.5),(f"#{sid}_h",0.3,0.4)]:
                anis.append(f'tl.from("{sel}",{{opacity:0,y:25,duration:{dd}}},{t+dy:.1f});')

        slides.append(f'<section id="{sid}" class="slide {layout}{" has-img" if has_img else ""}" data-start="{t:.1f}" data-duration="{dur:.1f}">{img_bg}<div class="slide-inner">{body}</div></section>')
        t += dur

    total = t
    css_vars = "\n".join(f"    {k}: {v};" for k,v in ds.items())
    title_tag = f"{storyboard[0].get('title','PPT')[:20]}"

    webgl_bg = ""
    if is_swiss:
        webgl_bg = """<canvas id="webgl-bg" style="position:fixed;top:0;left:0;z-index:0;opacity:0.03"></canvas>
<script>(function(){const c=document.getElementById('webgl-bg');if(!c)return;c.width=1920;c.height=1080;const x=c.getContext('2d');let t=0;function d(){x.clearRect(0,0,1920,1080);for(let i=0;i<5;i++){const cx=960+Math.sin(t*0.003+i)*400,cy=540+Math.cos(t*0.002+i)*300,r=200+Math.sin(t*0.001+i)*100;const g=x.createRadialGradient(cx,cy,0,cx,cy,r);g.addColorStop(0,'rgba(0,0,0,0.12)');g.addColorStop(1,'rgba(0,0,0,0)');x.beginPath();x.arc(cx,cy,r,0,Math.PI*2);x.fillStyle=g;x.fill()}t++;requestAnimationFrame(d)}d()})()</script>"""

    html = f'''<!DOCTYPE html><html lang="zh-CN">
<head><meta charset="UTF-8"><meta name="viewport" content="width={W},height={H}">
<script src="https://cdn.jsdelivr.net/npm/gsap@3.14.2/dist/gsap.min.js"></script><title>{title_tag}</title>
<style>
:root{{{css_vars}--w:{W}px;--h:{H}px;--margin:clamp(40px,6vw,120px);--fs-hero:clamp(72px,10vw,130px);--fs-h1:clamp(56px,7vw,100px);--fs-h2:clamp(42px,5vw,72px);--fs-body:clamp(22px,2.5vw,36px);--fs-small:clamp(16px,1.8vw,28px);--font-sans:"PingFang SC","Hiragino Sans GB","Microsoft YaHei",sans-serif;--font-display:"Georgia","Times New Roman",var(--font-sans),serif}}
*{{margin:0;padding:0;box-sizing:border-box}}
html,body{{width:var(--w);height:var(--h);overflow:hidden;background:var(--bg);font-family:var(--font-sans);color:var(--text);position:relative}}
.slide{{position:absolute;top:0;left:0;width:var(--w);height:var(--h);display:flex;align-items:center;justify-content:center;flex-direction:column;z-index:1;text-wrap:pretty}}
h1,h2,h3{{text-wrap:balance}}
.swiss-hero{{padding:0 var(--margin)}}.swiss-hero h1{{font-family:var(--font-display);font-size:var(--fs-hero);font-weight:900;line-height:1.08;letter-spacing:-0.5px;max-width:1500px}}
.swiss-hero .lead{{font-size:var(--fs-body);color:var(--textDim);margin-top:clamp(16px,2vw,32px);max-width:900px}}
.accent-bar{{position:absolute;top:0;left:0;width:1px;height:100%;background:var(--accent);transform-origin:left}}
.swiss-statement{{padding:0 var(--margin);text-align:left;width:1600px;position:relative}}
.swiss-statement .num{{font-size:clamp(140px,18vw,220px);font-weight:900;color:var(--accent);line-height:1;opacity:.1;position:absolute;top:60px;left:var(--margin)}}
.swiss-statement h2{{font-size:var(--fs-h1);font-weight:900;line-height:1.15;position:relative;z-index:1}}
.swiss-statement p{{font-size:var(--fs-body);color:var(--textDim);margin-top:20px;max-width:1000px}}
.swiss-split{{display:grid;grid-template-columns:1fr 1fr;gap:clamp(40px,5vw,80px);padding:0 var(--margin);width:100%}}
.swiss-split .col-l,.swiss-split .col-r{{display:flex;flex-direction:column;justify-content:center}}
.swiss-split h2{{font-size:var(--fs-h2);font-weight:900;line-height:1.15;margin-top:16px}}
.swiss-split p{{font-size:var(--fs-body);color:var(--textDim);line-height:1.5}}
.swiss-split .label{{font-size:var(--fs-small);font-weight:700;color:var(--accent);text-transform:uppercase;letter-spacing:4px}}
.swiss-quote{{padding:0 clamp(60px,8vw,160px);text-align:left;width:1600px}}
.swiss-quote .rule{{width:80px;height:1px;background:var(--accent);margin-bottom:40px}}
.swiss-quote blockquote{{font-family:var(--font-display);font-size:var(--fs-h2);font-weight:700;line-height:1.3;max-width:1400px}}
.swiss-quote cite{{font-size:var(--fs-small);color:var(--textDim);display:block;margin-top:32px}}
.swiss-kpi{{display:flex;flex-direction:column;align-items:center;gap:12px}}
.swiss-kpi .num{{font-size:clamp(120px,18vw,200px);font-weight:900;color:var(--accent);line-height:1}}
.swiss-kpi .title{{font-size:var(--fs-h2);font-weight:700}}
.swiss-kpi .desc{{font-size:var(--fs-small);color:var(--textDim)}}
h1{{font-family:var(--font-display);font-size:var(--fs-h1);font-weight:var(--fontTitle);line-height:1.15}}
h2{{font-size:var(--fs-h2);font-weight:var(--fontTitle);line-height:1.2}}
.subtitle{{font-size:var(--fs-body);color:var(--textDim);margin-top:20px;max-width:1000px}}
.tag{{display:inline-block;background:var(--accent);color:var(--bg);font-size:var(--fs-small);font-weight:800;padding:6px 22px;border-radius:2px;margin-bottom:24px;letter-spacing:2px;text-transform:uppercase}}
.img-bg{{position:absolute;top:0;left:0;width:100%;height:100%;background-size:cover;background-position:center;z-index:0;filter:brightness(0.55) saturate(0.8)}}
.img-overlay{{position:absolute;top:0;left:0;width:100%;height:100%;background:linear-gradient(180deg,rgba(0,0,0,0.3) 0%,rgba(0,0,0,0.65) 100%);z-index:0}}
.slide.has-img .slide-inner{{position:relative;z-index:1}}
.slide.has-img h1,.slide.has-img h2,.slide.has-img blockquote{{text-shadow:0 2px 20px rgba(0,0,0,0.5)}}
</style></head><body>
{webgl_bg}
<div id="comp" data-composition-id="main" data-start="0" data-duration="{total:.1f}" data-width="{W}" data-height="{H}">
{chr(10).join(audios)}{chr(10).join(slides)}
</div>
<script>window.__timelines=window.__timelines||{{}};const tl=gsap.timeline({{paused:true}});{chr(10).join(anis)}window.__timelines["main"]=tl</script>
</body></html>'''
    return html, total, img_to_copy

# ==================== P6: 快慢双渲染 ====================
def render_video(html_path, proj_dir, storyboard, voice_durs):
    """HyperFrames 渲染，失败则降级到 moviepy 快路径"""
    # 先试 HyperFrames
    env = {"PATH": os.path.expanduser("~/bin") + ":" + os.environ.get("PATH", "")}
    try:
        r = subprocess.run(["npx","hyperframes","render"], cwd=proj_dir, env=env,
                          capture_output=True, text=True, timeout=180)
        for line in (r.stdout + r.stderr).split("\n"):
            if line.strip().endswith(".mp4"):
                return line.strip().split()[-1]  # last token is mp4 path
    except: pass

    # 降级：moviepy 快路径（仅静态图+配音，无动画）
    print("  ⚡ HyperFrames 失败，降级 moviepy 快路径...")
    try:
        from moviepy import ImageClip, TextClip, ColorClip, AudioFileClip, concatenate_videoclips, CompositeVideoClip
        clips = []
        t_offset = 0
        for i, slide in enumerate(storyboard):
            dur = voice_durs[i] if i < len(voice_durs) else 2.0
            # 黑底+大字
            base = ColorClip((W, H), color=(0,0,0)).with_duration(dur)
            title = slide.get("title","")
            txt = TextClip(text=title, font=FONT_PATH, font_size=80, color="white",
                          stroke_color="black", stroke_width=3, size=(1600,None),
                          method="caption", text_align="center", duration=dur)
            txt = txt.with_position(("center","center"))
            clip = CompositeVideoClip([base, txt])
            voice_path = os.path.join(OUTPUT_DIR, f"voice_{i}.mp3")
            if os.path.exists(voice_path):
                try:
                    ac = AudioFileClip(voice_path)
                    clip = clip.with_audio(ac)
                except: pass
            clips.append(clip)
            t_offset += dur
        video = concatenate_videoclips(clips, method="compose")
        out = os.path.join(OUTPUT_DIR, "video_fast.mp4")
        video.write_videofile(out, fps=FPS, codec="libx264", bitrate=BITRATE, audio_codec="aac")
        return out
    except Exception as e:
        print(f"  ❌ 所有渲染路径失败: {e}")
        return None

# ==================== 主流程 ====================
def run():
    for d in [OUTPUT_DIR, os.path.join(OUTPUT_DIR,"images")]:
        os.makedirs(d, exist_ok=True)

    print("=" * 60)
    print("make_video.py v9 — 字幕+Few-shot+智能配色")
    print("=" * 60)

    # 1. 选题
    topic = ask_text("给我1个短视频爆款选题，15字以内，直接输出标题本身，不要编号")
    print(f"\n📌 选题: {topic}")

    # 2. 分镜（P2+P5）
    print("\n📝 生成分镜（JSON模式+自评重写）...")
    storyboard = generate_storyboard(topic)
    print(f"  ✅ {len(storyboard)} 页")

    # 3. 配音
    voice_durs = gen_voice(storyboard)

    # 4. 搜图（P1：多源+缓存）
    print(f"\n🖼 多源搜图（Pexels zh-CN → Unsplash → 缓存）...")
    images = {}
    for i, slide in enumerate(storyboard):
        kw = slide.get("highlight", "") or slide.get("title", "")[:10]
        if not kw: continue
        path, photog = fetch_image(kw, i)
        if path:
            images[i] = (path, photog)
            print(f"  [{i}] {kw[:15]} → {photog}")
        time.sleep(0.3)  # 请求间隔避免限流
    print(f"  命中 {len(images)}/{len(storyboard)}")

    # 5. 配色（AI实时生成 / 智能匹配 / 随机）
    theme_data, theme_name = pick_theme(topic)
    is_swiss = theme_data.get("--grid") == "swiss"
    print(f"\n🎨 主题: {theme_name} {'🇨🇭' if is_swiss else ''}")

    # 6. 生成 HTML
    html, total, img_to_copy = build_html(storyboard, theme_data, voice_durs, images)
    html_path = os.path.join(OUTPUT_DIR, "slides.html")
    with open(html_path, "w") as f: f.write(html)
    print(f"🌐 HTML: {html_path} ({total:.1f}s)")

    # 7. PPTX
    print("\n📄 导出 PPTX（图片+文字）...")
    pptx_path = export_pptx(storyboard, theme_data, images)

    # 8. HyperFrames 项目
    proj_dir = os.path.join(OUTPUT_DIR, "slides_proj")
    if not os.path.exists(proj_dir):
        os.makedirs(proj_dir)
        subprocess.run(["npx","hyperframes","init","slides_proj"], cwd=OUTPUT_DIR, capture_output=True, timeout=30)
    with open(os.path.join(proj_dir, "index.html"), "w") as f: f.write(html)
    for i in range(len(storyboard)):
        src = os.path.join(OUTPUT_DIR, f"voice_{i}.mp3")
        dst = os.path.join(proj_dir, f"voice_{i}.mp3")
        if os.path.exists(src): shutil.copy(src, dst)
    # 复制图片
    proj_img = os.path.join(proj_dir, "images")
    os.makedirs(proj_img, exist_ok=True)
    for img_path, idx in img_to_copy:
        dst = os.path.join(proj_img, f"img_{idx}.jpg")
        if os.path.exists(img_path): shutil.copy(img_path, dst)

    # 9. 渲染（P6：双路径）
    print(f"\n🎬 渲染视频...")
    video_path = render_video(html_path, proj_dir, storyboard, voice_durs)
    if video_path:
            # 10. BGM 混音
        print(f"\n🎵 BGM 混音...")
        bgm_src = BGM_PATH if os.path.exists(BGM_PATH) else None
        video_path = mix_bgm(video_path, bgm_src or "", video_path)

        # 11. 字幕（从分镜+配音时长生成SRT → 烧录）
        print(f"\n📝 字幕生成 + 烧录...")
        srt_path = generate_srt(storyboard, voice_durs)
        video_path = burn_subtitles(video_path, srt_path)

    # 完成
    print(f"\n{'='*60}")
    print(f"✅ 全部完成!")
    if video_path: print(f"  🎬 视频: {video_path}")
    if pptx_path: print(f"  📄 PPTX: {pptx_path}")
    print(f"  🌐 HTML: {html_path}")
    print(f"{'='*60}")

def run_from_storyboard(topic, storyboard):
    """跳过选题+分镜生成，直接用已有 storyboard 做视频"""
    for d in [OUTPUT_DIR, os.path.join(OUTPUT_DIR, "images")]:
        os.makedirs(d, exist_ok=True)
    print("=" * 60)
    print("make_video.py — 视频模式（从已有分镜）")
    print("=" * 60)
    print(f"\n📌 {topic} | {len(storyboard)} 页")

    voice_durs = gen_voice(storyboard)
    print(f"\n🖼 搜图...")
    images = {}
    for i, slide in enumerate(storyboard):
        kw = slide.get("highlight", "") or slide.get("title", "")[:10]
        if not kw: continue
        path, photog = fetch_image(kw, i)
        if path: images[i] = (path, photog); print(f"  [{i}] {kw[:12]} → {photog}")
        time.sleep(0.3)
    print(f"  命中 {len(images)}/{len(storyboard)}")
    theme_data, theme_name = pick_theme(topic)
    print(f"\n🎨 {theme_name}")
    html, total, img_to_copy = build_html(storyboard, theme_data, voice_durs, images)
    html_path = os.path.join(OUTPUT_DIR, "slides.html")
    with open(html_path, "w") as f: f.write(html)
    proj_dir = os.path.join(OUTPUT_DIR, "slides_proj")
    if not os.path.exists(proj_dir):
        os.makedirs(proj_dir)
        subprocess.run(["npx","hyperframes","init","slides_proj"], cwd=OUTPUT_DIR, capture_output=True, timeout=30)
    with open(os.path.join(proj_dir, "index.html"), "w") as f: f.write(html)
    for i in range(len(storyboard)):
        src = os.path.join(OUTPUT_DIR, f"voice_{i}.mp3")
        if os.path.exists(src): shutil.copy(src, os.path.join(proj_dir, f"voice_{i}.mp3"))
    proj_img = os.path.join(proj_dir, "images"); os.makedirs(proj_img, exist_ok=True)
    for img_path, idx in img_to_copy:
        if os.path.exists(img_path): shutil.copy(img_path, os.path.join(proj_img, f"img_{idx}.jpg"))
    print(f"\n🎬 渲染...")
    video_path = render_video(html_path, proj_dir, storyboard, voice_durs)
    if video_path:
        bgm_src = BGM_PATH if os.path.exists(BGM_PATH) else None
        video_path = mix_bgm(video_path, bgm_src or "", video_path)
        srt_path = generate_srt(storyboard, voice_durs)
        video_path = burn_subtitles(video_path, srt_path)
    print(f"\n✅ 完成! 🎬 {video_path}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "--storyboard" and len(sys.argv) > 2:
        # 从大纲文件取分镜，跳过选题+分镜生成
        with open(sys.argv[2]) as f:
            data = json.load(f)
        storyboard = data.get("storyboard", data.get("outline", {}).get("storyboard", data))
        if isinstance(storyboard, dict): storyboard = storyboard.get("storyboard", [])
        topic = data.get("topic", data.get("outline", {}).get("topic", "视频"))
        run_from_storyboard(topic, storyboard)
    else:
        run()
