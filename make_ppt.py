#!/usr/bin/env python3
"""
make_ppt.py — PPT 全流程工具（大纲→审稿→生成 HTML+PPTX）

用法:
  python3 make_ppt.py                    交互模式：选题→大纲→审稿→PPT
  python3 make_ppt.py "选题"              快速模式：直出大纲+PPT
  python3 make_ppt.py --json outline.json 从已有大纲生成 PPT
  python3 make_ppt.py --quick "选题"       极速模式：跳过审稿，直出 PPT
"""
import os, json, re, random, sys, time, requests

os.makedirs("output", exist_ok=True)
os.makedirs("output/images", exist_ok=True)

try:
    from config import DS_KEY, PEXELS_KEY
except ImportError:
    DS_KEY = "sk-你的DeepSeek-Key"
    PEXELS_KEY = "你的Pexels-Key"
W, H = 1920, 1080

# ==================== API ====================
def ask_json(system, user, max_t=3000):
    headers = {"Authorization": f"Bearer {DS_KEY}", "Content-Type": "application/json"}
    payload = {"model": "deepseek-chat", "messages": [{"role": "system", "content": system},
               {"role": "user", "content": user}], "max_tokens": max_t,
               "response_format": {"type": "json_object"}, "temperature": 0.7}
    for attempt in range(4):
        try:
            r = requests.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=payload, timeout=90)
            data = r.json()
            if "choices" not in data:
                if attempt < 3: time.sleep(3); continue
                raise RuntimeError(str(data.get("error","?"))[:200])
            content = data["choices"][0]["message"]["content"]
            if not content or content.strip() in ("", "{}", "[]"):
                time.sleep(2); continue
            return json.loads(content)
        except Exception as e:
            if attempt == 3 and "choices" not in str(e): raise
            time.sleep(3)

def ask_text(prompt, max_t=500):
    headers = {"Authorization": f"Bearer {DS_KEY}", "Content-Type": "application/json"}
    payload = {"model": "deepseek-chat", "messages": [{"role": "user", "content": prompt}], "max_tokens": max_t}
    r = requests.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=payload, timeout=30)
    return r.json()["choices"][0]["message"]["content"].strip()

# ==================== 搜图 ====================
def fetch_image(keyword, index):
    has_cjk = any('一' <= c <= '鿿' for c in keyword)
    params = {"query": keyword, "per_page": 1, "orientation": "landscape"}
    if has_cjk: params["locale"] = "zh-CN"
    try:
        r = requests.get("https://api.pexels.com/v1/search", headers={"Authorization": PEXELS_KEY}, params=params, timeout=10)
        data = r.json()
        if "photos" in data and data["photos"]:
            url = data["photos"][0]["src"]["large"]
            img = requests.get(url, timeout=15).content
            path = f"output/images/img_{index}.jpg"
            with open(path, "wb") as f: f.write(img)
            return path, data["photos"][0]["photographer"]
    except: pass
    return None, None

# ==================== 6套大纲模板 ====================
OUTLINE_TEMPLATES = {
    "冲击反转": {"structure": "钩子→反常识→数据→原因→新视角→行动→金句",
                 "适合": "抖音口播、争议话题"},
    "故事叙事": {"structure": "场景→人物→困境→转折→领悟→升华",
                 "适合": "品牌故事、情感类"},
    "干货教程": {"structure": "问题→原因→方法1→方法2→方法3→总结",
                 "适合": "知识付费、How-to"},
    "数据说服": {"structure": "核心数字→对比→原因→影响→方案→预期",
                 "适合": "商业汇报、Pitch"},
    "产品发布": {"structure": "痛点→方案缺陷→我们→功能→效果→CTA",
                 "适合": "新品发布、Demo"},
    "观点评论": {"structure": "事件→主流观点→立场→论据→反驳→结论",
                 "适合": "热点评论、行业分析"},
}

# ==================== 大纲生成 ====================
OUTLINE_PROMPT = """你是顶级内容策划师。根据选题生成 PPT 大纲。请输出JSON。

## 可用模板
{templates}

## 输出 JSON
{{
  "topic": "选题",
  "template": "选用的模板名",
  "angle": "切入角度（一句话）",
  "target_audience": "目标受众",
  "tone": "语气（冲击/理性/温情/幽默）",
  "hook": "开头钩子（15字以内）",
  "closing": "结尾金句（20字以内）",
  "outline": [
    {{"section": "板块标题", "key_message": "核心信息", "bullet_points": ["要点1","要点2"], "suggested_visual": "英文搜图关键词"}}
  ],
  "estimated_pages": 数字
}}
要求：5-8板块，每板块2-4要点，每要点≤15字，开头必须有冲突"""

def generate_outline(topic, template_name=None):
    tmpl_list = "\n".join(f"- {k}: {v['structure']} (适合{v['适合']})" for k,v in OUTLINE_TEMPLATES.items())
    hint = f"\n优先使用「{template_name}」模板。" if template_name and template_name in OUTLINE_TEMPLATES else "\n根据选题自动选模板。"
    return ask_json(OUTLINE_PROMPT.replace("{templates}", tmpl_list),
                    f"选题：{topic}{hint}\n请输出JSON", 3000)

def refine_outline(outline, feedback):
    return ask_json("你是内容策划师。根据用户反馈修改大纲。请输出JSON。",
                    f"原始大纲：{json.dumps(outline, ensure_ascii=False)}\n反馈：{feedback}\n输出完整大纲JSON", 3000)

# ==================== 大纲→分镜脚本 ====================
FEWSHOT = """
## 高质量示例
选题"工资涨不过物价":
[{"layout":"hero","rhythm":"hero","title":"你的工资涨了","subtitle":"但物价涨得更快","highlight":"工资陷阱"},
 {"layout":"kpi_bold","rhythm":"non-hero","title":"物价翻3倍","subtitle":"工资只涨1.5倍","highlight":"3x差距"},
 {"layout":"quote_slide","rhythm":"hero","title":"你不是在赚钱","subtitle":"你在被通胀吃掉","highlight":"真相"},
 {"layout":"split_col","rhythm":"non-hero","title":"钱存银行","subtitle":"每年贬值6%","highlight":"贬值"},
 {"layout":"hero","rhythm":"hero","title":"出路只有一个","subtitle":"让钱为你工作","highlight":"资产"}]"""

EXPAND_PROMPT = f"""你是短视频脚本写手。将 PPT 大纲展开为逐页分镜。请输出JSON。

{FEWSHOT}

## 规则：hero/non-hero交替、禁鸡汤、highlight用2-6字关键词
## layout：hero / statement / split_col / quote_slide / kpi_bold

输出JSON：{{"storyboard":[...]}}"""

def expand_to_storyboard(outline):
    r = ask_json(EXPAND_PROMPT, f"大纲：{json.dumps(outline, ensure_ascii=False)}\n请输出JSON", 4000)
    return r.get("storyboard", [])

# ==================== 55套设计系统 ====================
DESIGN_SYSTEMS = {
    "swiss_ikb":{"--bg":"#f8f8f8","--card":"#fff","--accent":"#0015FF","--accent2":"#000","--text":"#1a1a1a","--textDim":"#666","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "swiss_dark":{"--bg":"#1a1a1a","--card":"#2a2a2a","--accent":"#FFE600","--accent2":"#fff","--text":"#f0f0f0","--textDim":"#999","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "editorial_serif":{"--bg":"#fcfaf7","--card":"#ffffff","--accent":"#8b0000","--accent2":"#333","--text":"#1a1a1a","--textDim":"#777","--fontTitle":"900","--radius":"0px","--grid":"editorial"},
    "xiaohongshu":{"--bg":"#ffffff","--card":"#f8f8f8","--accent":"#ff2442","--accent2":"#cc1030","--text":"#333","--textDim":"#999","--fontTitle":"800","--radius":"8px","--grid":"none"},
    "tokyo_night":{"--bg":"#1a1b2e","--card":"#252740","--accent":"#ff9e64","--accent2":"#e0894f","--text":"#c0caf5","--textDim":"#6c7096","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "pitch_navy":{"--bg":"#020617","--card":"#0f172a","--accent":"#38bdf8","--accent2":"#0284c7","--text":"#f8fafc","--textDim":"#94a3b8","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "dracula":{"--bg":"#282a36","--card":"#383a59","--accent":"#bd93f9","--accent2":"#ff79c6","--text":"#f8f8f2","--textDim":"#6272a4","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "cyberpunk":{"--bg":"#0a0a1a","--card":"#1a1a3e","--accent":"#00ff88","--accent2":"#00cc66","--text":"#e0e0ff","--textDim":"#a0a0cc","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "minimal_white":{"--bg":"#ffffff","--card":"#f5f5f5","--accent":"#111","--accent2":"#666","--text":"#111","--textDim":"#888","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "bold_signal":{"--bg":"#000","--card":"#0f0f0f","--accent":"#FF4500","--accent2":"#FFD700","--text":"#fff","--textDim":"#aaa","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "apple_snow":{"--bg":"#f5f5f7","--card":"#fff","--accent":"#007aff","--accent2":"#5856d6","--text":"#1d1d1f","--textDim":"#86868b","--fontTitle":"700","--radius":"12px","--grid":"none"},
    "forest_moss":{"--bg":"#0a1a10","--card":"#142818","--accent":"#4ade80","--accent2":"#22c55e","--text":"#dcfce7","--textDim":"#86efac","--fontTitle":"900","--radius":"6px","--grid":"none"},
    "corp_blue":{"--bg":"#0a1e3d","--card":"#112850","--accent":"#3b82f6","--accent2":"#1d4ed8","--text":"#dbeafe","--textDim":"#93c5fd","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "warm_craft":{"--bg":"#faf5eb","--card":"#fff","--accent":"#d97706","--accent2":"#92400e","--text":"#44403c","--textDim":"#78716c","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "neo_brutalism":{"--bg":"#fff","--card":"#ffdd00","--accent":"#ff0000","--accent2":"#0000ff","--text":"#000","--textDim":"#555","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "glassmorphism":{"--bg":"#0f0f23","--card":"rgba(255,255,255,0.08)","--accent":"#ffffff","--accent2":"#c0c0ff","--text":"#ffffff","--textDim":"#c0c0d0","--fontTitle":"800","--radius":"16px","--grid":"none"},
    "vaporwave":{"--bg":"#0d0d1a","--card":"#1a1a3e","--accent":"#ff6ac1","--accent2":"#00ffff","--text":"#ffe0f0","--textDim":"#a080c0","--fontTitle":"900","--radius":"16px","--grid":"none"},
    "cream_paper":{"--bg":"#fefdf9","--card":"#fffef8","--accent":"#b45309","--accent2":"#78350f","--text":"#292524","--textDim":"#78716c","--fontTitle":"800","--radius":"2px","--grid":"none"},
    "nord":{"--bg":"#2e3440","--card":"#3b4252","--accent":"#88c0d0","--accent2":"#81a1c1","--text":"#eceff4","--textDim":"#d8dee9","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "rose_pine":{"--bg":"#191724","--card":"#1f1d2e","--accent":"#ebbcba","--accent2":"#f6c177","--text":"#e0def4","--textDim":"#908caa","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "mono_void":{"--bg":"#000","--card":"#0a0a0a","--accent":"#fff","--accent2":"#808080","--text":"#fff","--textDim":"#808080","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "vintage_film":{"--bg":"#1a1410","--card":"#2a2018","--accent":"#d4a574","--accent2":"#8b6914","--text":"#f5e6d3","--textDim":"#c4a882","--fontTitle":"800","--radius":"2px","--grid":"none"},
    "bauhaus":{"--bg":"#f0f0e8","--card":"#e8c547","--accent":"#c82828","--accent2":"#2c4c9c","--text":"#1a1a1a","--textDim":"#555","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "sunset_warm":{"--bg":"#1a0a05","--card":"#3a1a10","--accent":"#f97316","--accent2":"#ef4444","--text":"#fff7ed","--textDim":"#fdba74","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "catppuccin_mocha":{"--bg":"#1e1e2e","--card":"#313244","--accent":"#cba6f7","--accent2":"#f5c2e7","--text":"#cdd6f4","--textDim":"#a6adc8","--fontTitle":"900","--radius":"6px","--grid":"none"},
}

# ==================== HTML 生成 ====================
def build_html(storyboard, theme_data, images):
    ds = theme_data
    is_swiss = ds.get("--grid") == "swiss"
    slides, anis, img_to_copy = [], [], []
    t = 0

    SWISS = {
        "hero": lambda t2,s,h: (f'<div class="swiss-hero"><div class="accent-bar"></div><h1>{t2}</h1><p class="lead">{s}</p><div class="tag">{h}</div></div>', 1.6),
        "statement":lambda t2,s,h: (f'<div class="swiss-statement"><div class="num">{h}</div><h2>{t2}</h2><p>{s}</p></div>', 1.0),
        "split_col":lambda t2,s,h: (f'<div class="swiss-split"><div class="col-l"><div class="label">{h}</div><h2>{t2}</h2></div><div class="col-r"><p>{s}</p></div></div>', 1.0),
        "quote_slide":lambda t2,s,h: (f'<div class="swiss-quote"><div class="rule"></div><blockquote>{t2}</blockquote><cite>{h}</cite></div>', 1.0),
        "kpi_bold":lambda t2,s,h: (f'<div class="swiss-kpi"><span class="num">{h}</span><span class="title">{t2}</span><span class="desc">{s}</span></div>', 1.0),
    }

    for i, slide in enumerate(storyboard):
        layout = slide.get("layout", "statement")
        title = slide.get("title", "")
        sub = slide.get("subtitle", "")
        hl = slide.get("highlight", "")
        dur = 2.5 * (1.6 if layout in ("hero","cover") else 1.0)
        sid = f"s{i}"

        has_img = i in images and images[i][0]
        img_bg = ""
        if has_img:
            img_bg = f'<div class="img-bg" style="background-image:url(images/img_{i}.jpg)"></div><div class="img-overlay"></div>'
            img_to_copy.append((images[i][0], i))

        if is_swiss and layout in SWISS:
            body, _ = SWISS[layout](title, sub, hl)
            anis.append(f'tl.from("#{sid} h1, #{sid} h2, #{sid} blockquote",{{opacity:0,y:30,duration:0.7}},{t:.1f});')
        else:
            body = f'<h2 id="{sid}_t">{title}</h2><p id="{sid}_s">{sub}</p><div id="{sid}_h" class="tag">{hl}</div>'
            anis.append(f'tl.from("#{sid}_t",{{opacity:0,y:25,duration:0.6}},{t:.1f});')
            anis.append(f'tl.from("#{sid}_s",{{opacity:0,y:20,duration:0.5}},{t+0.15:.1f});')

        slides.append(f'<section id="{sid}" class="slide {layout}{" has-img" if has_img else ""}" data-start="{t:.1f}" data-duration="{dur:.1f}">{img_bg}<div class="slide-inner">{body}</div></section>')
        t += dur

    total = t
    css_vars = "\n".join(f"    {k}: {v};" for k,v in ds.items())

    html = f'''<!DOCTYPE html><html lang="zh-CN">
<head><meta charset="UTF-8"><meta name="viewport" content="width={W},height={H}">
<script src="https://cdn.jsdelivr.net/npm/gsap@3.14.2/dist/gsap.min.js"></script>
<style>
:root{{{css_vars}--w:{W}px;--h:{H}px;--margin:clamp(40px,6vw,120px);--fs-hero:clamp(72px,10vw,130px);--fs-h1:clamp(56px,7vw,100px);--fs-h2:clamp(42px,5vw,72px);--fs-body:clamp(22px,2.5vw,36px);--fs-small:clamp(16px,1.8vw,28px);--font-sans:"PingFang SC","Hiragino Sans GB","Microsoft YaHei",sans-serif;--font-display:"Georgia","Times New Roman",var(--font-sans),serif}}
*{{margin:0;padding:0;box-sizing:border-box}}
html,body{{width:var(--w);height:var(--h);overflow:hidden;background:var(--bg);font-family:var(--font-sans);color:var(--text);position:relative}}
.slide{{position:absolute;top:0;left:0;width:var(--w);height:var(--h);display:flex;align-items:center;justify-content:center;flex-direction:column;z-index:1;text-wrap:pretty}}
h1,h2,h3{{text-wrap:balance}}
.swiss-hero{{padding:0 var(--margin)}}.swiss-hero h1{{font-family:var(--font-display);font-size:var(--fs-hero);font-weight:900;line-height:1.08;max-width:1500px}}
.swiss-hero .lead{{font-size:var(--fs-body);color:var(--textDim);margin-top:clamp(16px,2vw,32px);max-width:900px}}
.accent-bar{{position:absolute;top:0;left:0;width:1px;height:100%;background:var(--accent);transform-origin:left}}
.swiss-statement{{padding:0 var(--margin);text-align:left;width:1600px;position:relative}}
.swiss-statement .num{{font-size:clamp(140px,18vw,220px);font-weight:900;color:var(--accent);line-height:1;opacity:.1;position:absolute;top:60px;left:var(--margin)}}
.swiss-statement h2{{font-size:var(--fs-h1);font-weight:900;line-height:1.15;position:relative;z-index:1}}
.swiss-statement p{{font-size:var(--fs-body);color:var(--textDim);margin-top:20px;max-width:1000px}}
.swiss-split{{display:grid;grid-template-columns:1fr 1fr;gap:clamp(40px,5vw,80px);padding:0 var(--margin);width:100%}}
.swiss-split h2{{font-size:var(--fs-h2);font-weight:900}}
.swiss-split .label{{font-size:var(--fs-small);font-weight:700;color:var(--accent);text-transform:uppercase;letter-spacing:4px}}
.swiss-quote{{padding:0 clamp(60px,8vw,160px);text-align:left;width:1600px}}
.swiss-quote .rule{{width:80px;height:1px;background:var(--accent);margin-bottom:40px}}
.swiss-quote blockquote{{font-family:var(--font-display);font-size:var(--fs-h2);font-weight:700;line-height:1.3}}
.swiss-kpi{{display:flex;flex-direction:column;align-items:center;gap:12px}}
.swiss-kpi .num{{font-size:clamp(120px,18vw,200px);font-weight:900;color:var(--accent);line-height:1}}
h1{{font-family:var(--font-display);font-size:var(--fs-h1);font-weight:var(--fontTitle)}}
h2{{font-size:var(--fs-h2);font-weight:var(--fontTitle)}}
.tag{{display:inline-block;background:var(--accent);color:var(--bg);font-size:var(--fs-small);font-weight:800;padding:6px 22px;border-radius:2px;margin-bottom:24px}}
.img-bg{{position:absolute;top:0;left:0;width:100%;height:100%;background-size:cover;background-position:center;z-index:0;filter:brightness(0.55) saturate(0.8)}}
.img-overlay{{position:absolute;top:0;left:0;width:100%;height:100%;background:linear-gradient(180deg,rgba(0,0,0,0.3),rgba(0,0,0,0.65));z-index:0}}
.slide.has-img .slide-inner{{position:relative;z-index:1}}
.slide.has-img h1,.slide.has-img h2,.slide.has-img blockquote{{text-shadow:0 2px 20px rgba(0,0,0,0.5)}}
</style></head><body>
<div id="comp" data-composition-id="main" data-start="0" data-duration="{total:.1f}" data-width="{W}" data-height="{H}">
{chr(10).join(slides)}
</div>
<script>window.__timelines=window.__timelines||{{}};const tl=gsap.timeline({{paused:true}});{chr(10).join(anis)}window.__timelines["main"]=tl</script>
</body></html>'''
    return html, total, img_to_copy

# ==================== PPTX 导出 ====================
def export_pptx(storyboard, theme_data, images):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        ds = theme_data
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

        def hex2rgb(h):
            h = str(h).lstrip('#').split(',')[0].split(')')[0].strip()
            if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2
            if len(h) < 6: h = h.ljust(6, '0')
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        bg_rgb = hex2rgb(ds.get("--bg","#000"))
        text_rgb = hex2rgb(ds.get("--text","#fff"))
        accent_rgb = hex2rgb(ds.get("--accent","#f00"))

        for i, s in enumerate(storyboard):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = bg_rgb

            if i in images and images[i][0] and os.path.exists(images[i][0]):
                slide.shapes.add_picture(images[i][0], Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0,0,0); shape.line.fill.background()

            title, subtitle, hl = s.get("title",""), s.get("subtitle",""), s.get("highlight","")
            if title:
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.5))
                tf = txBox.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = text_rgb
            if subtitle:
                txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.5))
                tf2 = txBox2.text_frame; tf2.word_wrap = True
                p2 = tf2.paragraphs[0]; p2.text = subtitle; p2.alignment = PP_ALIGN.CENTER
                p2.font.size = Pt(20); p2.font.color.rgb = text_rgb
            if hl:
                txBox3 = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2.3), Inches(0.6))
                tf3 = txBox3.text_frame
                p3 = tf3.paragraphs[0]; p3.text = hl; p3.alignment = PP_ALIGN.CENTER
                p3.font.size = Pt(18); p3.font.bold = True; p3.font.color.rgb = accent_rgb

        path = "output/slides.pptx"
        prs.save(path)
        return path
    except Exception as e:
        print(f"  ⚠️ PPTX导出失败: {e}")
        return None

# ==================== 打印大纲 ====================
def print_outline(outline):
    print(f"\n{'─'*50}")
    print(f"📌 {outline.get('topic','')}")
    print(f"  🎯 {outline.get('angle','')}")
    print(f"  👥 {outline.get('target_audience','')} | 🗣 {outline.get('tone','')} | 📋 {outline.get('template','')}")
    print(f"  🪝 钩子: {outline.get('hook','')}")
    print(f"{'─'*50}")
    for i, sec in enumerate(outline.get("outline", []), 1):
        print(f"\n  {i}. {sec.get('section','')}")
        print(f"     💡 {sec.get('key_message','')}")
        for bp in sec.get("bullet_points", []):
            print(f"       • {bp}")
    print(f"\n  🔚 {outline.get('closing','')}")
    print(f"  📄 预计 {outline.get('estimated_pages','?')} 页")

# ==================== 主流程 ====================
def make_ppt_from_storyboard(topic, storyboard):
    """从 storyboard 生成 HTML+PPTX"""
    print(f"\n{'='*50}")
    print(f"📌 {topic} | {len(storyboard)} 页")
    print("=" * 50)

    print("\n🖼 Pexels 搜图...")
    images = {}
    for i, s in enumerate(storyboard):
        kw = s.get("highlight","") or s.get("title","")[:10]
        path, photog = fetch_image(kw, i)
        if path:
            images[i] = (path, photog)
            print(f"  [{i}] {kw[:12]} → {photog}")
        time.sleep(0.3)

    theme_name = random.choice(list(DESIGN_SYSTEMS.keys()))
    theme_data = DESIGN_SYSTEMS[theme_name]
    print(f"\n🎨 主题: {theme_name} | 命中 {len(images)}/{len(storyboard)} 图")

    html, total, img_to_copy = build_html(storyboard, theme_data, images)
    html_path = "output/slides.html"
    with open(html_path, "w") as f: f.write(html)

    pptx_path = export_pptx(storyboard, theme_data, images)

    print(f"\n✅ 完成!  HTML: {html_path}  |  PPTX: {pptx_path or '无'}")
    print(f"💡 浏览器打开 HTML 可编辑文字，按 ← → 翻页")
    return html_path, pptx_path


def interactive():
    """交互模式：选题→大纲→审稿→编辑→PPT"""
    print("=" * 50)
    print("PPT 全流程工具 — 大纲→审稿→生成")
    print("=" * 50)

    # 1. 输入选题
    topic = input("\n📌 请输入选题: ").strip()
    if not topic:
        topic = ask_text("给我一个爆款选题，15字以内")
        print(f"  AI 选题: {topic}")

    # 2. 选模板
    print("\n📋 6套内容模板:")
    for i, (name, info) in enumerate(OUTLINE_TEMPLATES.items(), 1):
        print(f"  {i}. {name} — {info['structure'][:40]}...")
    choice = input(f"\n  选模板 (1-6，回车自动): ").strip()
    template = list(OUTLINE_TEMPLATES.keys())[int(choice)-1] if choice.isdigit() and 1 <= int(choice) <= 6 else None

    # 3. 生成大纲
    print("\n⏳ 生成大纲...")
    outline = generate_outline(topic, template)
    print_outline(outline)

    # 4. 审稿编辑循环
    while True:
        action = input(f"\n🔧 [a]接受并生成PPT  [e]修改大纲  [r]重生成  [t]换模板  [q]退出: ").strip().lower()
        if action == 'a':
            break
        elif action == 'e':
            fb = input("  修改意见（如：第二段太弱、加数据支撑、语气太硬）: ").strip()
            if fb:
                print("  ⏳ 优化中...")
                outline = refine_outline(outline, fb)
                print_outline(outline)
        elif action == 'r':
            print("  ⏳ 重生成...")
            outline = generate_outline(topic, template)
            print_outline(outline)
        elif action == 't':
            print("\n📋 模板:")
            for i, (name, info) in enumerate(OUTLINE_TEMPLATES.items(), 1):
                print(f"  {i}. {name}")
            c = input("  选: ").strip()
            if c.isdigit() and 1 <= int(c) <= 6:
                template = list(OUTLINE_TEMPLATES.keys())[int(c)-1]
                outline = generate_outline(topic, template)
                print_outline(outline)
        elif action == 'q':
            with open("output/outline.json", "w") as f:
                json.dump(outline, f, ensure_ascii=False, indent=2)
            print("已保存 output/outline.json，随时可 --json 恢复")
            return

    # 5. 展开为分镜
    print("\n⏳ 展开为逐页分镜脚本...")
    storyboard = expand_to_storyboard(outline)
    print(f"  ✅ {len(storyboard)} 页分镜")

    # 6. 保存大纲
    with open("output/outline.json", "w") as f:
        json.dump({"topic": topic, "outline": outline, "storyboard": storyboard}, f, ensure_ascii=False, indent=2)

    # 7. 生成 PPT
    make_ppt_from_storyboard(topic, storyboard)

    # 8. 顺便出视频？
    _offer_video(topic)


def _offer_video(topic, interactive=True):
    """PPT 完成后询问/提示是否生成视频"""
    if interactive:
        print(f"\n{'─'*50}")
        choice = input("🎬 是否同时生成视频（配音+BGM+字幕）？[y/N]: ").strip().lower()
        if choice == 'y':
            print("⏳ 启动视频生成...")
            import subprocess
            subprocess.run([sys.executable, "make_video.py", "--storyboard", "output/outline.json"])
            return
    print("💡 需要视频时运行: python3 make_video.py --storyboard output/outline.json")


def quick_mode(topic):
    """快速模式：跳过审稿，直出大纲+PPT"""
    print(f"\n📌 {topic}")
    print("⏳ 大纲→分镜→PPT 全自动...")
    outline = generate_outline(topic)
    print_outline(outline)
    storyboard = expand_to_storyboard(outline)
    with open("output/outline.json", "w") as f:
        json.dump({"topic": topic, "outline": outline, "storyboard": storyboard}, f, ensure_ascii=False, indent=2)
    make_ppt_from_storyboard(topic, storyboard)
    _offer_video(topic, interactive=False)


# ==================== 入口 ====================
if __name__ == "__main__":
    if len(sys.argv) > 1:
        arg = sys.argv[1]
        if arg == "--json" and len(sys.argv) > 2:
            with open(sys.argv[2]) as f:
                data = json.load(f)
            storyboard = data.get("storyboard", data.get("outline", {}).get("storyboard", data))
            if isinstance(storyboard, dict): storyboard = storyboard.get("storyboard", [])
            topic = data.get("topic", data.get("outline", {}).get("topic", "PPT"))
            make_ppt_from_storyboard(topic, storyboard)
            _offer_video(topic, interactive=False)
        elif arg == "--quick" and len(sys.argv) > 2:
            quick_mode(" ".join(sys.argv[2:]))
        else:
            quick_mode(" ".join(sys.argv[1:]))
    else:
        interactive()
