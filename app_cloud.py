"""云端版 — 仅 PPT 生成，轻量，永久在线"""
import os, json, re, random, time, requests
from flask import Flask, request, render_template_string, send_file

app = Flask(__name__)
try: from config import DS_KEY, PEXELS_KEY
except: DS_KEY = os.environ.get("DS_KEY",""); PEXELS_KEY = os.environ.get("PEXELS_KEY","")

W, H = 1920, 1080
OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(f"{OUTPUT_DIR}/images", exist_ok=True)

# ========== DeepSeek API ==========
def ask_json(system, user, max_t=3000):
    for attempt in range(3):
        try:
            r = requests.post("https://api.deepseek.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {DS_KEY}","Content-Type":"application/json"},
                json={"model":"deepseek-chat","messages":[{"role":"system","content":system},
                      {"role":"user","content":user}],"max_tokens":max_t,"response_format":{"type":"json_object"}},
                timeout=90)
            d = r.json()
            return json.loads(d["choices"][0]["message"]["content"])
        except: time.sleep(2)
    return {}

def ask_text(prompt, max_t=500):
    r = requests.post("https://api.deepseek.com/v1/chat/completions",
        headers={"Authorization": f"Bearer {DS_KEY}","Content-Type":"application/json"},
        json={"model":"deepseek-chat","messages":[{"role":"user","content":prompt}],"max_tokens":max_t}, timeout=30)
    return r.json()["choices"][0]["message"]["content"].strip()

# ========== Pexels ==========
def fetch_image(keyword, index):
    has_cjk = any('一' <= c <= '鿿' for c in keyword)
    params = {"query": keyword, "per_page": 1, "orientation": "landscape"}
    if has_cjk: params["locale"] = "zh-CN"
    try:
        r = requests.get("https://api.pexels.com/v1/search", headers={"Authorization": PEXELS_KEY}, params=params, timeout=10)
        data = r.json()
        if "photos" in data and data["photos"]:
            url = data["photos"][0]["src"]["large"]
            img = requests.get(url, timeout=15).content
            path = f"{OUTPUT_DIR}/images/img_{index}.jpg"
            with open(path, "wb") as f: f.write(img)
            return path
    except: pass
    return None

# ========== 分镜生成 ==========
STORYBOARD_PROMPT = """你是短视频导演。根据选题生成逐页分镜。请输出JSON。
## 规则：hero/non-hero交替、禁鸡汤、highlight用2-6字关键词、每句≤15字
## layout：hero / statement / split_col / quote_slide / kpi_bold
## 示例 — 选题"工资涨不过物价":
[{"layout":"hero","rhythm":"hero","title":"你的工资涨了","subtitle":"但物价涨得更快","highlight":"工资陷阱"},
 {"layout":"kpi_bold","rhythm":"non-hero","title":"物价翻3倍","subtitle":"工资只涨1.5倍","highlight":"3x差距"},
 {"layout":"quote_slide","rhythm":"hero","title":"你不是在赚钱","subtitle":"你在被通胀吃掉","highlight":"真相"}]
输出JSON：{"storyboard":[...]}"""

def gen_storyboard(topic):
    for _ in range(2):
        try:
            r = ask_json(STORYBOARD_PROMPT, f"选题：{topic}\n请输出JSON", 3000)
            sb = r.get("storyboard", [])
            if isinstance(sb, list) and len(sb) >= 3: return sb
        except: continue
    return [{"layout":"hero","rhythm":"hero","title":topic,"subtitle":"","highlight":"核心"}]

# ========== 25套主题 ==========
THEMES = {
    "swiss_dark":{"--bg":"#1a1a1a","--card":"#2a2a2a","--accent":"#FFE600","--accent2":"#fff","--text":"#f0f0f0","--textDim":"#999","--fontTitle":"900","--radius":"0px","--grid":"swiss"},
    "editorial":{"--bg":"#fcfaf7","--card":"#fff","--accent":"#8b0000","--accent2":"#333","--text":"#1a1a1a","--textDim":"#777","--fontTitle":"900","--radius":"0px","--grid":"editorial"},
    "cyberpunk":{"--bg":"#0a0a1a","--card":"#1a1a3e","--accent":"#00ff88","--accent2":"#00cc66","--text":"#e0e0ff","--textDim":"#a0a0cc","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "minimal":{"--bg":"#fff","--card":"#f5f5f5","--accent":"#111","--accent2":"#666","--text":"#111","--textDim":"#888","--fontTitle":"700","--radius":"4px","--grid":"none"},
    "bold":{"--bg":"#000","--card":"#0f0f0f","--accent":"#FF4500","--accent2":"#FFD700","--text":"#fff","--textDim":"#aaa","--fontTitle":"900","--radius":"0px","--grid":"none"},
    "tokyo":{"--bg":"#1a1b2e","--card":"#252740","--accent":"#ff9e64","--accent2":"#e0894f","--text":"#c0caf5","--textDim":"#6c7096","--fontTitle":"900","--radius":"8px","--grid":"none"},
    "warm":{"--bg":"#faf5eb","--card":"#fff","--accent":"#d97706","--accent2":"#92400e","--text":"#44403c","--textDim":"#78716c","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "forest":{"--bg":"#0a1a10","--card":"#142818","--accent":"#4ade80","--accent2":"#22c55e","--text":"#dcfce7","--textDim":"#86efac","--fontTitle":"900","--radius":"6px","--grid":"none"},
    "nord":{"--bg":"#2e3440","--card":"#3b4252","--accent":"#88c0d0","--accent2":"#81a1c1","--text":"#eceff4","--textDim":"#d8dee9","--fontTitle":"800","--radius":"4px","--grid":"none"},
    "sunset":{"--bg":"#1a0a05","--card":"#3a1a10","--accent":"#f97316","--accent2":"#ef4444","--text":"#fff7ed","--textDim":"#fdba74","--fontTitle":"900","--radius":"8px","--grid":"none"},
}

# ========== HTML 生成 ==========
def build_html(storyboard, theme_data, images):
    ds = theme_data
    is_swiss = ds.get("--grid") == "swiss"
    slides, anis = [], []
    t = 0
    SWISS = {
        "hero": lambda t2,s,h: (f'<div class="swiss-hero"><div class="accent-bar"></div><h1>{t2}</h1><p class="lead">{s}</p><div class="tag">{h}</div></div>',),
        "statement":lambda t2,s,h: (f'<div class="swiss-statement"><div class="num">{h}</div><h2>{t2}</h2><p>{s}</p></div>',),
        "split_col":lambda t2,s,h: (f'<div class="swiss-split"><div class="col-l"><div class="label">{h}</div><h2>{t2}</h2></div><div class="col-r"><p>{s}</p></div></div>',),
        "quote_slide":lambda t2,s,h: (f'<div class="swiss-quote"><div class="rule"></div><blockquote>{t2}</blockquote><cite>{h}</cite></div>',),
        "kpi_bold":lambda t2,s,h: (f'<div class="swiss-kpi"><span class="num">{h}</span><span class="title">{t2}</span><span class="desc">{s}</span></div>',),
    }
    for i, slide in enumerate(storyboard):
        layout = slide.get("layout","statement")
        title = slide.get("title",""); sub = slide.get("subtitle",""); hl = slide.get("highlight","")
        dur = 2.5 * (1.6 if layout in ("hero","cover") else 1.0)
        sid = f"s{i}"
        has_img = i in images
        img_bg = f'<div class="img-bg" style="background-image:url(images/img_{i}.jpg)"></div><div class="img-overlay"></div>' if has_img else ""
        if is_swiss and layout in SWISS:
            body, _ = SWISS[layout](title, sub, hl)
        else:
            body = f'<h2 id="{sid}_t">{title}</h2><p id="{sid}_s">{sub}</p><div id="{sid}_h" class="tag">{hl}</div>'
        anis.append(f'tl.from("#{sid} h1, #{sid} h2, #{sid} blockquote, #{sid}_t",{{opacity:0,y:30,duration:0.7}},{t:.1f});')
        slides.append(f'<section id="{sid}" class="slide {layout}{" has-img" if has_img else ""}" data-start="{t:.1f}" data-duration="{dur:.1f}">{img_bg}<div class="slide-inner">{body}</div></section>')
        t += dur
    total = t
    css_vars = "\n".join(f"    {k}: {v};" for k,v in ds.items())
    return f'''<!DOCTYPE html><html lang="zh-CN">
<head><meta charset="UTF-8"><meta name="viewport" content="width={W},height={H}">
<script src="https://cdn.jsdelivr.net/npm/gsap@3.14.2/dist/gsap.min.js"></script>
<style>
:root{{{css_vars}--w:{W}px;--h:{H}px;--margin:clamp(40px,6vw,120px);--fs-hero:clamp(72px,10vw,130px);--fs-h1:clamp(56px,7vw,100px);--fs-h2:clamp(42px,5vw,72px);--fs-body:clamp(22px,2.5vw,36px);--font-sans:"PingFang SC","Hiragino Sans GB",sans-serif;--font-display:"Georgia",serif}}
*{{margin:0;padding:0;box-sizing:border-box}}
html,body{{width:var(--w);height:var(--h);overflow:hidden;background:var(--bg);font-family:var(--font-sans);color:var(--text)}}
.slide{{position:absolute;top:0;left:0;width:var(--w);height:var(--h);display:flex;align-items:center;justify-content:center;flex-direction:column;z-index:1;text-wrap:pretty}}
.swiss-hero h1{{font-family:var(--font-display);font-size:var(--fs-hero);font-weight:900;line-height:1.08}}
.swiss-hero .lead{{font-size:var(--fs-body);color:var(--textDim);margin-top:24px}}
.accent-bar{{position:absolute;top:0;left:0;width:1px;height:100%;background:var(--accent)}}
.swiss-statement{{text-align:left;width:1600px;position:relative}}
.swiss-statement .num{{font-size:200px;font-weight:900;color:var(--accent);line-height:1;opacity:.1;position:absolute;top:60px;left:80px}}
.swiss-statement h2{{font-size:var(--fs-h1);font-weight:900;position:relative;z-index:1}}
.swiss-split{{display:grid;grid-template-columns:1fr 1fr;gap:60px;width:100%;padding:0 80px}}
.swiss-split h2{{font-size:var(--fs-h2);font-weight:900}}
.swiss-quote{{text-align:left;width:1600px;padding:0 120px}}
.swiss-quote .rule{{width:80px;height:1px;background:var(--accent);margin-bottom:40px}}
.swiss-quote blockquote{{font-family:var(--font-display);font-size:var(--fs-h2);font-weight:700}}
.swiss-kpi{{display:flex;flex-direction:column;align-items:center;gap:12px}}
.swiss-kpi .num{{font-size:180px;font-weight:900;color:var(--accent);line-height:1}}
h1{{font-family:var(--font-display);font-size:var(--fs-h1);font-weight:var(--fontTitle)}}
h2{{font-size:var(--fs-h2);font-weight:var(--fontTitle)}}
.tag{{display:inline-block;background:var(--accent);color:var(--bg);font-size:22px;font-weight:800;padding:6px 22px;border-radius:2px;margin-bottom:24px}}
.img-bg{{position:absolute;top:0;left:0;width:100%;height:100%;background-size:cover;background-position:center;z-index:0;filter:brightness(0.55)}}
.img-overlay{{position:absolute;top:0;left:0;width:100%;height:100%;background:linear-gradient(180deg,rgba(0,0,0,0.3),rgba(0,0,0,0.65));z-index:0}}
.slide.has-img .slide-inner{{position:relative;z-index:1;text-shadow:0 2px 20px rgba(0,0,0,0.5)}}
</style></head><body>
<div data-composition-id="main" data-start="0" data-duration="{total:.1f}" data-width="{W}" data-height="{H}">
{chr(10).join(slides)}
</div>
<script>window.__timelines={{}};const tl=gsap.timeline({{paused:true}});{chr(10).join(anis)}window.__timelines["main"]=tl</script>
</body></html>'''

# ========== Web UI ==========
HTML_PAGE = '''<!DOCTYPE html><html lang="zh-CN">
<head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>AI PPT</title>
<style>
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:-apple-system,sans-serif;background:#0a0a15;color:#fff;min-height:100vh;padding:20px}
h1{font-size:24px;text-align:center;margin:16px 0;background:linear-gradient(135deg,#00ff88,#00ccff);-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.card{background:#1a1a2e;border-radius:12px;padding:20px;margin-bottom:16px}
input{width:100%;padding:14px;border-radius:10px;border:1px solid #333;background:#0f0f23;color:#fff;font-size:16px;margin-bottom:10px}
button{width:100%;padding:14px;border-radius:10px;border:none;font-size:18px;font-weight:700;background:linear-gradient(135deg,#00ff88,#00cc66);color:#000;cursor:pointer}
.log{background:#0a0a16;border-radius:10px;padding:14px;font-size:13px;color:#aaa;max-height:200px;overflow-y:auto;white-space:pre-wrap;margin-top:10px}
.result{margin-top:12px}.result a{display:block;padding:8px;color:#00ff88;text-decoration:none}
</style></head><body>
<h1>🎯 AI PPT 工作室</h1>
<div class="card">
  <input id="topic" placeholder="📌 输入选题，如：为什么越省钱越穷" autofocus>
  <button onclick="gen()">⚡ 生成 PPT</button>
  <div class="log" id="log"></div>
  <div id="res"></div>
</div>
<script>
async function gen(){
  const t=document.getElementById('topic').value.trim();if(!t)return;
  document.getElementById('log').textContent='⏳ 生成中...';
  const r=await fetch('/generate',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({topic:t})});
  const d=await r.json();
  document.getElementById('log').textContent=d.log||'完成';
  if(d.files){let h='';d.files.forEach(f=>{h+=`<a href="/download/${f.name}" download>📥 ${f.label} (${f.size})</a>`});document.getElementById('res').innerHTML=h}
}
</script></body></html>'''

@app.route('/')
def index(): return render_template_string(HTML_PAGE)

@app.route('/generate', methods=['POST'])
def generate():
    topic = request.json.get('topic','').strip()
    try:
        storyboard = gen_storyboard(topic)
        images = {}
        for i, s in enumerate(storyboard):
            kw = s.get("highlight","") or s.get("title","")[:10]
            path = fetch_image(kw, i)
            if path: images[i] = path
            time.sleep(0.3)
        theme = random.choice(list(THEMES.keys()))
        html = build_html(storyboard, THEMES[theme], images)
        html_path = f"{OUTPUT_DIR}/slides.html"
        with open(html_path, "w") as f: f.write(html)
        # PPTX
        pptx_path = f"{OUTPUT_DIR}/slides.pptx"
        try:
            from pptx import Presentation; from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor; from pptx.enum.text import PP_ALIGN
            ds = THEMES[theme]
            def h2r(h):
                h=str(h).lstrip('#').split(',')[0].split(')')[0].strip()
                if len(h)==3: h=h[0]*2+h[1]*2+h[2]*2
                return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
            prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
            for i,s in enumerate(storyboard):
                slide=prs.slides.add_slide(prs.slide_layouts[6])
                bg=slide.background;fill=bg.fill;fill.solid();fill.fore_color.rgb=h2r(ds.get("--bg","#000"))
                title=s.get("title","");sub=s.get("subtitle","");hl=s.get("highlight","")
                if title:
                    tb=slide.shapes.add_textbox(Inches(0.8),Inches(2.5),Inches(11.5),Inches(2.5))
                    tf=tb.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=title;p.alignment=PP_ALIGN.CENTER
                    p.font.size=Pt(44);p.font.bold=True;p.font.color.rgb=h2r(ds.get("--text","#fff"))
                if sub:
                    tb2=slide.shapes.add_textbox(Inches(0.8),Inches(4.8),Inches(11.5),Inches(1.5))
                    tf2=tb2.text_frame;tf2.word_wrap=True;p2=tf2.paragraphs[0];p2.text=sub;p2.alignment=PP_ALIGN.CENTER
                    p2.font.size=Pt(20);p2.font.color.rgb=h2r(ds.get("--textDim","#999"))
            prs.save(pptx_path)
        except: pptx_path = None

        files = [{"name":"slides.html","label":"🌐 HTML 幻灯片","size":"OK"}]
        if pptx_path: files.append({"name":"slides.pptx","label":"📄 PPTX","size":"OK"})
        return {"log": f"✅ {topic}\n{len(storyboard)} 页 | 主题: {theme} | {len(images)} 张图", "files": files}
    except Exception as e:
        return {"log": f"❌ {e}", "files": []}

@app.route('/download/<path:filename>')
def download(filename):
    path = os.path.join(OUTPUT_DIR, filename)
    if os.path.exists(path): return send_file(path, as_attachment=True)
    return "Not found", 404

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 7890)))
